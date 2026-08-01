from copy import deepcopy
from io import BytesIO
from math import sqrt
from pathlib import Path
import re

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Pt


MAX_SOURCE_SLIDES = 100
MAX_OUTPUT_SLIDES = 50
MAX_ADDED_SLIDES = 20
MAX_REMOVED_SLIDES = 3
MAX_INVENTORY_ELEMENT_CHARS = 2500
EMU_PER_POINT = 12700
MIN_BODY_FONT_PT = 11
MIN_TITLE_FONT_PT = 18
PLACEHOLDER_PATTERN = re.compile(r"\[[^\[\]\n]{1,100}\]")


def _clean_text(value):
    return " ".join(str(value or "").split())


def _font_size_points(text_frame, fallback):
    """텍스트 상자에 지정된 첫 글자 크기를 읽고 테마 글꼴이면 기본값을 씁니다."""

    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size:
                return round(run.font.size.pt, 1)
        if paragraph.font.size:
            return round(paragraph.font.size.pt, 1)
    return fallback


def _estimate_text_capacity(width, height, font_size_pt, current_text, is_title):
    """상자 크기와 글자 크기로 겹침 없이 들어갈 권장 글자 수를 계산합니다."""

    width_pt = max(1, width / EMU_PER_POINT)
    height_pt = max(1, height / EMU_PER_POINT)
    chars_per_line = max(4, int(width_pt / max(font_size_pt * 0.9, 1)))
    line_count = max(1, int(height_pt / max(font_size_pt * 1.25, 1)))
    estimated = int(chars_per_line * line_count * 0.85)
    current_length = len(_clean_text(current_text))
    capacity = max(current_length, estimated)
    if is_title:
        return min(max(capacity, 18), 60)
    return min(max(capacity, 35), MAX_INVENTORY_ELEMENT_CHARS)


def _element_metadata(target, text, text_frame, width, height, kind, is_title=False):
    fallback_size = 28 if is_title else 16
    font_size_pt = _font_size_points(text_frame, fallback_size)
    return {
        "target": target,
        "text": text,
        "kind": "title" if is_title else kind,
        "font_size_pt": font_size_pt,
        "max_chars": _estimate_text_capacity(
            width,
            height,
            font_size_pt,
            text,
            is_title,
        ),
    }


def _slide_elements(slide):
    """AI와 개정 엔진이 함께 사용할 슬라이드의 텍스트 위치를 만듭니다."""

    elements = []
    title_shape_id = (
        slide.shapes.title.shape_id
        if slide.shapes.title is not None
        else None
    )
    for shape_index, shape in enumerate(slide.shapes):
        if getattr(shape, "has_table", False):
            for row_index, row in enumerate(shape.table.rows):
                for cell_index, cell in enumerate(row.cells):
                    text = cell.text.strip()
                    if text:
                        elements.append(
                            _element_metadata(
                                target=(
                                    f"shape-{shape_index}-cell-{row_index}-{cell_index}"
                                ),
                                text=text,
                                text_frame=cell.text_frame,
                                width=shape.table.columns[cell_index].width,
                                height=row.height,
                                kind="table_cell",
                            )
                        )
            continue

        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if text:
                elements.append(
                    _element_metadata(
                        target=f"shape-{shape_index}",
                        text=text,
                        text_frame=shape.text_frame,
                        width=shape.width,
                        height=shape.height,
                        kind="text",
                        is_title=shape.shape_id == title_shape_id,
                    )
                )

    return elements


def _classify_slide(slide_number, title, elements):
    """슬라이드의 쓰임을 분류해 AI가 알맞은 레이아웃을 고르게 합니다."""

    text_length = sum(len(_clean_text(element["text"])) for element in elements)
    combined_text = f"{title} {' '.join(element['text'] for element in elements)}".lower()
    if slide_number == 1:
        return "cover"
    if "contents" in combined_text or "목차" in combined_text:
        return "contents"
    if len(elements) <= 1 and text_length <= 160:
        return "divider"
    if sum(element["kind"] == "table_cell" for element in elements) >= 4:
        return "matrix"
    if any(keyword in combined_text for keyword in ("일정", "schedule", "마일스톤")):
        return "timeline"
    if any(keyword in combined_text for keyword in ("조직", "인력", "organization")):
        return "organization"
    if any(keyword in combined_text for keyword in ("위험", "품질", "보안", "risk")):
        return "governance"
    if any(keyword in combined_text for keyword in ("실적", "역량", "인증", "evidence")):
        return "evidence"
    if any(keyword in combined_text for keyword in ("요약", "summary", "핵심 제안")):
        return "summary"
    if any(keyword in combined_text for keyword in ("절차", "방법", "수행", "process")):
        return "process"
    return "content"


def extract_pptx_inventory(source_path, max_slides=MAX_SOURCE_SLIDES):
    """원본 PPTX의 슬라이드 번호와 수정 가능한 텍스트 위치를 읽습니다."""

    source_path = Path(source_path)
    if source_path.suffix.lower() != ".pptx":
        raise ValueError("기존 PowerPoint 제안서(.pptx)만 개정할 수 있습니다.")

    presentation = Presentation(str(source_path))
    if len(presentation.slides) > max_slides:
        raise ValueError(
            f"입찰 제안서는 최대 {max_slides}장까지 처리할 수 있습니다."
        )

    inventory = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        elements = _slide_elements(slide)
        title = ""
        if slide.shapes.title is not None:
            title = slide.shapes.title.text.strip()
        if not title and elements:
            title = elements[0]["text"].splitlines()[0]

        inventory.append(
            {
                "slide_number": slide_number,
                "title": title or f"{slide_number}페이지",
                "elements": elements,
                "role": _classify_slide(slide_number, title, elements),
            }
        )

    if not inventory:
        raise ValueError("원본 제안서에 슬라이드가 없습니다.")
    return inventory


def build_inventory_context(inventory, max_chars=50000):
    """슬라이드 구조를 AI가 읽기 쉬운 짧은 텍스트로 변환합니다."""

    lines = []
    per_slide_limit = max(220, max_chars // max(1, len(inventory)))
    for slide in inventory:
        slide_lines = [
            (
                f"[슬라이드 {slide['slide_number']}] "
                f"역할={slide.get('role', 'content')} | {slide['title']}"
            ),
        ]
        for element in slide["elements"]:
            text = _clean_text(element["text"])[:MAX_INVENTORY_ELEMENT_CHARS]
            slide_lines.append(
                f"- {element['target']} | 종류={element.get('kind', 'text')} | "
                f"글자={element.get('font_size_pt', 16)}pt | "
                f"권장 최대={element.get('max_chars', 200)}자: {text}"
            )

        # 긴 제안서도 마지막 페이지까지 빠짐없이 전달하도록 페이지별 길이를 나눕니다.
        lines.append("\n".join(slide_lines)[:per_slide_limit])

    return "\n\n".join(lines)


def get_content_template_numbers(inventory):
    """표지와 간지를 제외한 추가 슬라이드용 본문 템플릿 번호를 반환합니다."""

    return [
        slide["slide_number"]
        for slide in inventory
        if slide.get("role") not in {"cover", "contents", "divider"}
        and len(slide.get("elements", [])) >= 2
    ]


def _replace_text_frame(text_frame, revised_text, font_size_pt=None):
    """첫 글자의 서식을 유지하며 텍스트만 교체합니다."""

    text_frame.word_wrap = True
    paragraphs = text_frame.paragraphs
    first_paragraph = paragraphs[0]
    if first_paragraph.runs:
        first_paragraph.runs[0].text = revised_text
        if font_size_pt:
            first_paragraph.runs[0].font.size = Pt(font_size_pt)
        for run in first_paragraph.runs[1:]:
            run.text = ""
    else:
        first_paragraph.text = revised_text
        if font_size_pt:
            first_paragraph.font.size = Pt(font_size_pt)

    for paragraph in paragraphs[1:]:
        for run in paragraph.runs:
            run.text = ""


def _target_text_frame(slide, target):
    parts = target.split("-")
    if len(parts) < 2 or parts[0] != "shape":
        return None

    try:
        shape = slide.shapes[int(parts[1])]
    except (IndexError, ValueError):
        return None

    if len(parts) == 2 and getattr(shape, "has_text_frame", False):
        return shape.text_frame

    if (
        len(parts) == 5
        and parts[2] == "cell"
        and getattr(shape, "has_table", False)
    ):
        try:
            return shape.table.cell(int(parts[3]), int(parts[4])).text_frame
        except (IndexError, ValueError):
            return None

    return None


def _apply_text_changes(slide, text_changes):
    applied = []
    warnings = []
    element_map = {
        element["target"]: element
        for element in _slide_elements(slide)
    }

    for change in text_changes:
        target = change.get("target", "")
        text_frame = _target_text_frame(slide, target)
        if text_frame is None:
            warnings.append(f"{target}: 텍스트 위치를 찾지 못했습니다.")
            continue

        original_text = change.get("original_text", "")
        current_text = text_frame.text
        clean_original = _clean_text(original_text)
        clean_current = _clean_text(current_text)
        if (
            clean_original
            and clean_original not in clean_current
            and clean_current not in clean_original
        ):
            warnings.append(f"{target}: 원문이 달라 자동 수정을 건너뛰었습니다.")
            continue

        revised_text = str(change.get("revised_text", "")).strip()
        element = element_map.get(target, {})
        max_chars = int(element.get("max_chars", MAX_INVENTORY_ELEMENT_CHARS))
        font_size_pt = float(element.get("font_size_pt", 16))
        revised_length = len(_clean_text(revised_text))
        adjusted_font_size = font_size_pt

        if revised_length > max_chars:
            adjusted_font_size = font_size_pt * sqrt(max_chars / revised_length)
            minimum_size = (
                MIN_TITLE_FONT_PT
                if element.get("kind") == "title"
                else MIN_BODY_FONT_PT
            )
            if adjusted_font_size < minimum_size:
                warnings.append(
                    f"{target}: 권장 {max_chars}자를 크게 넘어 자동 수정을 건너뛰었습니다."
                )
                continue
            adjusted_font_size = max(minimum_size, adjusted_font_size)

        _replace_text_frame(
            text_frame,
            revised_text,
            font_size_pt=adjusted_font_size,
        )
        applied.append(
            {
                "target": target,
                "content_label": (
                    change.get("content_label")
                    or _clean_text(current_text)[:30]
                    or "본문"
                ),
                "summary": change.get("reason", "") or "새 공고 기준으로 수정",
                "char_count": revised_length,
                "max_chars": max_chars,
                "font_size_pt": round(adjusted_font_size, 1),
            }
        )

    return applied, warnings


def _clone_slide(presentation, source_slide):
    """원본 슬라이드의 도형과 연결된 이미지를 복사해 같은 디자인을 재사용합니다."""

    cloned_slide = presentation.slides.add_slide(source_slide.slide_layout)
    for shape in list(cloned_slide.shapes):
        cloned_slide.shapes._spTree.remove(shape.element)

    relationship_ids = {}
    for relationship_id, relationship in source_slide.part.rels.items():
        if relationship.reltype in {RT.SLIDE_LAYOUT, RT.NOTES_SLIDE}:
            continue

        if relationship.is_external:
            new_relationship_id = cloned_slide.part.rels.get_or_add_ext_rel(
                relationship.reltype,
                relationship.target_ref,
            )
        else:
            new_relationship_id = cloned_slide.part.rels.get_or_add(
                relationship.reltype,
                relationship.target_part,
            )
        relationship_ids[relationship_id] = new_relationship_id

    for shape in source_slide.shapes:
        shape_element = deepcopy(shape.element)
        for element in shape_element.iter():
            for attribute, value in list(element.attrib.items()):
                if value in relationship_ids:
                    element.set(attribute, relationship_ids[value])
        cloned_slide.shapes._spTree.insert_element_before(
            shape_element,
            "p:extLst",
        )

    return cloned_slide


def _append_revision_note(slide, message):
    """변경 메모를 화면이 아닌 PowerPoint 발표자 노트에 남깁니다."""

    text_frame = slide.notes_slide.notes_text_frame
    if text_frame is None:
        return

    existing_note = text_frame.text.strip()
    note = f"[AI 제안서 개정 메모]\n{message.strip()}"
    text_frame.text = f"{existing_note}\n\n{note}".strip()


def _slide_id_element(presentation, slide):
    for slide_id in presentation.slides._sldIdLst:
        if presentation.part.related_slide(slide_id.rId).part is slide.part:
            return slide_id
    raise ValueError("슬라이드 순서를 확인할 수 없습니다.")


def _set_slide_order(presentation, ordered_slides):
    slide_id_list = presentation.slides._sldIdLst
    ordered_part_ids = {id(slide.part) for slide in ordered_slides}

    for slide_id in list(slide_id_list):
        slide = presentation.part.related_slide(slide_id.rId)
        if id(slide.part) not in ordered_part_ids:
            presentation.part.drop_rel(slide_id.rId)
            slide_id_list.remove(slide_id)

    ordered_ids = [
        _slide_id_element(presentation, slide)
        for slide in ordered_slides
    ]
    for slide_id in ordered_ids:
        slide_id_list.remove(slide_id)
        slide_id_list.append(slide_id)


def _discard_slide(presentation, slide):
    """검증에 실패한 임시 복제 슬라이드를 프레젠테이션에서 제거합니다."""

    slide_id = _slide_id_element(presentation, slide)
    presentation.part.drop_rel(slide_id.rId)
    presentation.slides._sldIdLst.remove(slide_id)


def inspect_proposal_quality(file_bytes, revision_log=None):
    """생성본의 미교체 문구, 빈 슬라이드와 적용 실패를 자동 점검합니다."""

    presentation = Presentation(BytesIO(file_bytes))
    unresolved_placeholders = []
    empty_slide_numbers = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        texts = [
            _clean_text(element.get("text", ""))
            for element in _slide_elements(slide)
            if _clean_text(element.get("text", ""))
        ]
        if len("".join(texts)) < 20:
            empty_slide_numbers.append(slide_number)

        for text in texts:
            for marker in PLACEHOLDER_PATTERN.findall(text):
                if re.fullmatch(
                    r"\[(?:출처|근거)\s*\d+(?:\s*,\s*\d+)*\]",
                    marker,
                ):
                    continue
                item = {"slide_number": slide_number, "marker": marker}
                if item not in unresolved_placeholders:
                    unresolved_placeholders.append(item)

    apply_warnings = []
    for entry in revision_log or []:
        for warning in entry.get("warnings", []):
            apply_warnings.append(
                {
                    "slide_number": entry.get("output_slide_number"),
                    "message": warning,
                }
            )

    review_items = []
    if unresolved_placeholders:
        slides = sorted(
            {item["slide_number"] for item in unresolved_placeholders}
        )
        review_items.append(
            "자리표시자가 남은 슬라이드를 확인해 주세요: "
            + ", ".join(map(str, slides))
        )
    if empty_slide_numbers:
        review_items.append(
            "내용이 거의 없는 슬라이드를 확인해 주세요: "
            + ", ".join(map(str, empty_slide_numbers))
        )
    if apply_warnings:
        review_items.append(
            f"텍스트 적용 경고 {len(apply_warnings)}건을 확인해 주세요."
        )

    return {
        "passed": not review_items,
        "unresolved_placeholders": unresolved_placeholders,
        "empty_slide_numbers": empty_slide_numbers,
        "apply_warnings": apply_warnings,
        "review_items": review_items,
    }


def build_proposal_pptx(
    source_path,
    bid_notice,
    revision_plan,
    max_source_slides=MAX_SOURCE_SLIDES,
    max_output_slides=MAX_OUTPUT_SLIDES,
):
    """Bid2 템플릿에 AI 작성 계획을 적용해 새 PPTX를 만듭니다."""

    inventory = extract_pptx_inventory(source_path, max_slides=max_source_slides)
    presentation = Presentation(str(source_path))
    original_slides = list(presentation.slides)
    source_slide_count = len(original_slides)
    if source_slide_count == 0:
        raise ValueError("원본 제안서에 슬라이드가 없습니다.")

    requested_changes = revision_plan.get("slide_changes", [])
    removed_numbers = set()
    for change in requested_changes:
        slide_number = int(change.get("slide_number", 0))
        if (
            change.get("action") == "REMOVE"
            and 1 < slide_number <= source_slide_count
            and len(removed_numbers) < MAX_REMOVED_SLIDES
        ):
            removed_numbers.add(slide_number)

    revision_log = []
    slide_log_entries = {}
    additions_by_position = {}
    added_slide_entries = []

    allowed_templates = set(get_content_template_numbers(inventory))
    available_additions = min(
        MAX_ADDED_SLIDES,
        max(0, max_output_slides - (source_slide_count - len(removed_numbers))),
    )

    # 추가 슬라이드는 사용자의 본문 디자인만 복제하며 표지와 간지는 차단합니다.
    for addition in revision_plan.get("added_slides", []):
        if len(added_slide_entries) >= available_additions:
            break

        template_number = int(addition.get("template_slide_number", 0))
        text_changes = addition.get("text_changes", [])
        if (
            template_number not in allowed_templates
            or len(text_changes) < 2
        ):
            continue

        after_number = int(addition.get("after_slide_number", source_slide_count))
        after_number = min(max(after_number, 0), source_slide_count)
        new_slide = _clone_slide(
            presentation,
            original_slides[template_number - 1],
        )
        applied, warnings = _apply_text_changes(new_slide, text_changes)
        if len(applied) < 2 or warnings:
            _discard_slide(presentation, new_slide)
            continue

        title = addition.get("title") or "추가 제안 내용"
        reason = addition.get("reason", "")
        _append_revision_note(new_slide, f"추가 슬라이드\n{title}\n{reason}")
        additions_by_position.setdefault(after_number, []).append(new_slide)

        entry = {
            "source_slide_number": None,
            "output_slide_number": None,
            "action": "ADD",
            "title": title,
            "reason": reason,
            "changes": applied,
            "warnings": [],
        }
        revision_log.append(entry)
        added_slide_entries.append((new_slide, entry))

    for change in requested_changes:
        slide_number = int(change.get("slide_number", 0))
        if not 1 <= slide_number <= source_slide_count:
            continue

        action = change.get("action", "REVIEW")
        slide = original_slides[slide_number - 1]
        title = change.get("title") or f"{slide_number}페이지"
        reason = change.get("reason", "")
        applied = []
        warnings = []

        if action == "REMOVE":
            if slide_number not in removed_numbers:
                action = "REVIEW"
                warnings.append(
                    "원본 보존 규칙에 따라 자동 삭제하지 않고 검토 대상으로 남겼습니다."
                )
                _append_revision_note(
                    slide,
                    f"{slide_number}페이지 삭제 여부 검토\n{reason}",
                )
        elif action == "UPDATE":
            applied, warnings = _apply_text_changes(
                slide,
                change.get("text_changes", []),
            )
            _append_revision_note(
                slide,
                f"{slide_number}페이지 수정\n{reason}",
            )
        else:
            _append_revision_note(
                slide,
                f"{slide_number}페이지 담당자 검토 필요\n{reason}",
            )

        entry = {
            "source_slide_number": slide_number,
            "output_slide_number": None,
            "action": action,
            "title": title,
            "reason": reason,
            "changes": applied,
            "warnings": warnings,
        }
        revision_log.append(entry)
        slide_log_entries[id(slide.part)] = entry

    ordered_slides = list(additions_by_position.get(0, []))
    for slide_number, slide in enumerate(original_slides, start=1):
        if slide_number not in removed_numbers:
            ordered_slides.append(slide)
        ordered_slides.extend(additions_by_position.get(slide_number, []))

    if not ordered_slides:
        raise ValueError("모든 슬라이드를 삭제할 수는 없습니다.")

    _set_slide_order(presentation, ordered_slides)
    output_numbers = {
        id(slide.part): index
        for index, slide in enumerate(ordered_slides, start=1)
    }
    for slide_part_id, entry in slide_log_entries.items():
        entry["output_slide_number"] = output_numbers.get(slide_part_id)
    for slide, entry in added_slide_entries:
        entry["output_slide_number"] = output_numbers.get(id(slide.part))

    buffer = BytesIO()
    presentation.save(buffer)
    file_bytes = buffer.getvalue()
    quality_review = inspect_proposal_quality(file_bytes, revision_log)
    return {
        "file_bytes": file_bytes,
        "filename": f"revised-proposal-{bid_notice.bid_ntce_no}.pptx",
        "content_type": (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ),
        "revision_log": revision_log,
        "source_slide_count": source_slide_count,
        "output_slide_count": len(ordered_slides),
        "quality_review": quality_review,
    }

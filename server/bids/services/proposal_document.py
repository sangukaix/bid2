from datetime import date
from io import BytesIO
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Pt as PptxPt


def clear_docx_body(document):
    """원본 Word의 스타일과 여백은 유지하고 이전 제안 내용만 제거합니다."""

    body = document._element.body
    for child in list(body):
        if not child.tag.endswith("sectPr"):
            body.remove(child)


def add_docx_bullet(document, text):
    try:
        paragraph = document.add_paragraph(style="List Bullet")
    except KeyError:
        paragraph = document.add_paragraph()
    paragraph.add_run(text)


def build_docx(source_path, bid_notice, strategy, draft):
    source_path = Path(source_path)
    if source_path.suffix.lower() == ".docx":
        document = Document(str(source_path))
        clear_docx_body(document)
    else:
        document = Document()

    normal_style = document.styles["Normal"]
    normal_style.font.name = normal_style.font.name or "Malgun Gothic"
    normal_style.font.size = Pt(10.5)

    title = document.add_heading(draft["proposal_title"], level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle = document.add_paragraph(draft["subtitle"])
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta = document.add_paragraph(
        f"공고번호 {bid_notice.bid_ntce_no}  |  작성일 {date.today().isoformat()}"
    )
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    document.add_page_break()

    document.add_heading("제안 요약", level=1)
    document.add_paragraph(draft["executive_summary"])

    document.add_heading("수주 전략", level=1)
    for theme in strategy.get("win_themes", []):
        add_docx_bullet(document, theme)
    document.add_page_break()

    remaining_content_pages = max(0, draft.get("page_limit", 30) - 4)
    for section in draft.get("sections", []):
        for page_index, page in enumerate(section.get("pages", [])):
            if remaining_content_pages <= 0:
                break

            if page_index == 0:
                document.add_heading(section["title"], level=1)
                if section.get("purpose"):
                    purpose = document.add_paragraph()
                    purpose.add_run("작성 목적: ").bold = True
                    purpose.add_run(section["purpose"])

            document.add_heading(page["title"], level=2)
            document.add_paragraph(page["content"][:900])
            for point in page.get("key_points", [])[:5]:
                add_docx_bullet(document, point)

            if page_index == 0 and section.get("company_evidence"):
                document.add_heading("회사 근거", level=2)
                for evidence in section["company_evidence"][:5]:
                    add_docx_bullet(document, evidence)

            if page_index == 0 and section.get("source_numbers"):
                document.add_paragraph(
                    "공고 근거: "
                    + ", ".join(
                        f"[출처 {number}]"
                        for number in section["source_numbers"]
                    )
                )

            remaining_content_pages -= 1
            if remaining_content_pages > 0:
                document.add_page_break()

        if remaining_content_pages <= 0:
            break

    document.add_page_break()
    document.add_heading("최종 확인 사항", level=1)
    for item in draft.get("final_checklist", []):
        add_docx_bullet(document, f"□ {item}")

    sources = draft.get("sources", [])
    if sources:
        document.add_heading("참고한 공고 문서", level=1)
        for source in sources:
            add_docx_bullet(
                document,
                f"[출처 {source['number']}] {source['file_name']} · {source['location']}",
            )

    buffer = BytesIO()
    document.save(buffer)
    return {
        "file_bytes": buffer.getvalue(),
        "filename": f"proposal-{bid_notice.bid_ntce_no}.docx",
        "content_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    }


def remove_all_slides(presentation):
    """PowerPoint의 테마와 슬라이드 마스터는 유지하고 기존 슬라이드를 제거합니다."""

    slide_ids = presentation.slides._sldIdLst
    for slide_id in list(slide_ids):
        presentation.part.drop_rel(slide_id.rId)
        slide_ids.remove(slide_id)


def add_pptx_slide(presentation, title, bullets):
    layout_index = 1 if len(presentation.slide_layouts) > 1 else 0
    slide = presentation.slides.add_slide(presentation.slide_layouts[layout_index])
    if slide.shapes.title:
        slide.shapes.title.text = title

    body_shape = next(
        (shape for shape in slide.placeholders if shape != slide.shapes.title and hasattr(shape, "text_frame")),
        None,
    )
    if body_shape is None:
        return

    text_frame = body_shape.text_frame
    text_frame.clear()
    for index, bullet in enumerate(bullets[:7]):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = PptxPt(18)


def build_pptx(source_path, bid_notice, strategy, draft):
    source_path = Path(source_path)
    if source_path.suffix.lower() == ".pptx":
        presentation = Presentation(str(source_path))
        remove_all_slides(presentation)
    else:
        presentation = Presentation()

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    if title_slide.shapes.title:
        title_slide.shapes.title.text = draft["proposal_title"]
    subtitle_shape = next(
        (shape for shape in title_slide.placeholders if shape != title_slide.shapes.title and hasattr(shape, "text_frame")),
        None,
    )
    if subtitle_shape:
        subtitle_shape.text = (
            f"{draft['subtitle']}\n공고번호 {bid_notice.bid_ntce_no}\n{date.today().isoformat()}"
        )

    add_pptx_slide(presentation, "제안 요약", [draft["executive_summary"]])
    add_pptx_slide(presentation, "핵심 수주 전략", strategy.get("win_themes", []))

    remaining_content_pages = max(0, draft.get("page_limit", 30) - 4)
    for section in draft.get("sections", []):
        for page in section.get("pages", []):
            if remaining_content_pages <= 0:
                break

            bullets = page.get("key_points") or [page.get("content", "")[:900]]
            add_pptx_slide(
                presentation,
                f"{section['title']} · {page['title']}",
                bullets,
            )
            remaining_content_pages -= 1

        if remaining_content_pages <= 0:
            break

    add_pptx_slide(
        presentation,
        "최종 확인 사항",
        [f"□ {item}" for item in draft.get("final_checklist", [])],
    )

    buffer = BytesIO()
    presentation.save(buffer)
    return {
        "file_bytes": buffer.getvalue(),
        "filename": f"proposal-{bid_notice.bid_ntce_no}.pptx",
        "content_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }


def build_proposal_file(source_path, output_format, bid_notice, strategy, draft):
    if output_format == "docx":
        return build_docx(source_path, bid_notice, strategy, draft)
    if output_format == "pptx":
        return build_pptx(source_path, bid_notice, strategy, draft)
    raise ValueError("지원하지 않는 제안서 출력 형식입니다.")

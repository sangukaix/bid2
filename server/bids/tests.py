from datetime import datetime, timedelta
from pathlib import Path
import tempfile
from unittest.mock import patch
import zipfile
from zoneinfo import ZoneInfo

from django.contrib.auth import get_user_model
from django.core.files.uploadedfile import SimpleUploadedFile
from django.db import IntegrityError, transaction
from django.test import TestCase
from django.test import override_settings
from django.utils import timezone
from rest_framework.authtoken.models import Token

from bids.management.commands.sync_bids import (
    determine_deadline_status,
    update_latest_notices,
)
from bids.models import (
    BidAnalysis,
    BidChatMessage,
    BidNotice,
    BidProposal,
    CompanyDocument,
    CompanyKnowledgeItem,
    CompanyProfile,
    RecommendedBid,
    SavedBid,
)
from bids.services.g2b_api import fetch_bid_notices
from bids.services.rag.extract_document import extract_document
from bids.services.rag.retriever import search_bid_documents
from bids.services.recommendation import match_user_recommendations


class SignupViewTests(TestCase):
    def setUp(self):
        self.signup_data = {
            "username": "new-user",
            "email": "new@example.com",
            "password": "testpass123",
            "password_confirm": "testpass123",
        }

    def test_회원가입으로_사용자를_생성한다(self):
        response = self.client.post(
            "/api/auth/signup/",
            self.signup_data,
            content_type="application/json",
        )

        user = get_user_model().objects.get(username="new-user")
        self.assertEqual(response.status_code, 201)
        self.assertEqual(response.json()["user"]["email"], "new@example.com")
        self.assertTrue(user.check_password("testpass123"))

    def test_4자리_숫자_비밀번호로_회원가입할_수_있다(self):
        response = self.client.post(
            "/api/auth/signup/",
            {
                "username": "short-password-user",
                "email": "short@example.com",
                "password": "1234",
                "password_confirm": "1234",
            },
            content_type="application/json",
        )

        self.assertEqual(response.status_code, 201)

    def test_같은_아이디로_다시_가입할_수_없다(self):
        get_user_model().objects.create_user(
            username="new-user",
            email="first@example.com",
            password="testpass123",
        )

        response = self.client.post(
            "/api/auth/signup/",
            self.signup_data,
            content_type="application/json",
        )

        self.assertEqual(response.status_code, 400)
        self.assertIn("username", response.json())


class LoginViewTests(TestCase):
    def setUp(self):
        self.user = get_user_model().objects.create_user(
            username="login-user",
            email="login@example.com",
            password="testpass123",
        )

    def test_로그인하면_token을_발급한다(self):
        response = self.client.post(
            "/api/auth/login/",
            {"username": "login-user", "password": "testpass123"},
            content_type="application/json",
        )

        self.assertEqual(response.status_code, 200)
        self.assertTrue(response.json()["token"])
        self.assertEqual(response.json()["user"]["username"], "login-user")
        self.assertNotIn("password", response.json())

    def test_비밀번호가_틀리면_로그인할_수_없다(self):
        response = self.client.post(
            "/api/auth/login/",
            {"username": "login-user", "password": "wrong-password"},
            content_type="application/json",
        )

        self.assertEqual(response.status_code, 400)
        self.assertIn("non_field_errors", response.json())

    def test_로그아웃하면_token을_삭제한다(self):
        token = Token.objects.create(user=self.user)

        response = self.client.post(
            "/api/auth/logout/",
            HTTP_AUTHORIZATION=f"Token {token.key}",
        )

        self.assertEqual(response.status_code, 200)
        self.assertFalse(Token.objects.filter(user=self.user).exists())

    def test_로그인하지_않으면_로그아웃할_수_없다(self):
        response = self.client.post("/api/auth/logout/")

        self.assertEqual(response.status_code, 401)


class CompanyProfileViewTests(TestCase):
    def setUp(self):
        self.user = get_user_model().objects.create_user(
            username="profile-user",
            password="test-password",
        )
        self.profile_data = {
            "company_name": "테스트회사",
            "business_registration_number": "000-00-00000",
            "representative_name": "김대표",
            "address": "서울특별시",
            "industry": "소프트웨어 개발",
            "related_industries": "정보통신, 연구·컨설팅",
            "main_business": "웹 서비스 개발",
            "preferred_keywords": "홈페이지 구축",
        }

    def test_로그인하지_않으면_회사_프로필에_접근할_수_없다(self):
        response = self.client.get("/api/company-profile/")

        self.assertEqual(response.status_code, 401)

    def test_회사_프로필이_없으면_null을_반환한다(self):
        self.client.force_login(self.user)

        response = self.client.get("/api/company-profile/")

        self.assertEqual(response.status_code, 200)
        self.assertIsNone(response.json()["profile"])

    def test_회사_프로필을_최초_저장한다(self):
        self.client.force_login(self.user)

        response = self.client.post(
            "/api/company-profile/",
            self.profile_data,
            content_type="application/json",
        )

        self.assertEqual(response.status_code, 201)
        self.assertEqual(CompanyProfile.objects.get().user, self.user)
        self.assertEqual(response.json()["profile"]["company_name"], "테스트회사")
        self.assertEqual(
            response.json()["profile"]["related_industries"],
            "정보통신, 연구·컨설팅",
        )

    def test_기존_회사_프로필의_일부를_수정한다(self):
        self.client.force_login(self.user)
        CompanyProfile.objects.create(user=self.user, **self.profile_data)

        response = self.client.patch(
            "/api/company-profile/",
            {"company_name": "수정된 회사"},
            content_type="application/json",
        )

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()["profile"]["company_name"], "수정된 회사")
        self.assertEqual(CompanyProfile.objects.get().industry, "소프트웨어 개발")

    def test_관련_업종은_최대_4개까지_저장한다(self):
        self.client.force_login(self.user)
        self.profile_data["related_industries"] = "교육, IT, 제조, 건설, 의료"

        response = self.client.post(
            "/api/company-profile/",
            self.profile_data,
            content_type="application/json",
        )

        self.assertEqual(response.status_code, 400)
        self.assertIn("related_industries", response.json())


class BusinessRegistrationViewTests(TestCase):
    def setUp(self):
        self.user = get_user_model().objects.create_user(
            username="registration-user",
            password="test-password",
        )

    def test_로그인하지_않으면_사업자등록증을_분석할_수_없다(self):
        response = self.client.post(
            "/api/company-profile/business-registration/",
            {"file": SimpleUploadedFile("registration.png", b"image", content_type="image/png")},
        )

        self.assertEqual(response.status_code, 401)

    @patch("bids.views.extract_business_registration")
    def test_사업자등록증에서_확인된_기본정보만_반환한다(self, extract_mock):
        self.client.force_login(self.user)
        extract_mock.return_value = {
            "company_name": "테스트회사",
            "business_registration_number": "123-45-67890",
            "representative_name": "김대표",
            "address": "서울특별시 강남구",
        }

        response = self.client.post(
            "/api/company-profile/business-registration/",
            {"file": SimpleUploadedFile("registration.png", b"image", content_type="image/png")},
        )

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()["company_name"], "테스트회사")
        self.assertEqual(response.json()["address"], "서울특별시 강남구")
        extract_mock.assert_called_once()

    def test_파일이_없으면_사업자등록증을_분석하지_않는다(self):
        self.client.force_login(self.user)

        response = self.client.post("/api/company-profile/business-registration/")

        self.assertEqual(response.status_code, 400)
        self.assertIn("error", response.json())


class CompanyDocumentViewTests(TestCase):
    def setUp(self):
        self.user = get_user_model().objects.create_user(
            username="document-user",
            password="test-password",
        )
        self.media_directory = tempfile.TemporaryDirectory()
        self.settings_override = override_settings(MEDIA_ROOT=self.media_directory.name)
        self.settings_override.enable()
        self.addCleanup(self.settings_override.disable)
        self.addCleanup(self.media_directory.cleanup)

    def test_로그인한_사용자가_제안서를_업로드하고_조회한다(self):
        self.client.force_login(self.user)
        uploaded_file = SimpleUploadedFile(
            "sample-proposal.docx",
            b"sample proposal",
            content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )

        response = self.client.post(
            "/api/company-documents/",
            {
                "file": uploaded_file,
                "document_type": "proposal",
            },
        )
        list_response = self.client.get("/api/company-documents/")

        self.assertEqual(response.status_code, 201)
        self.assertEqual(list_response.status_code, 200)
        self.assertEqual(list_response.json()["count"], 1)
        self.assertEqual(list_response.json()["items"][0]["original_name"], "sample-proposal.docx")

    def test_PDF_회사_문서는_업로드할_수_없다(self):
        self.client.force_login(self.user)

        response = self.client.post(
            "/api/company-documents/",
            {
                "file": SimpleUploadedFile(
                    "sample.pdf",
                    b"pdf",
                    content_type="application/pdf",
                ),
                "document_type": "proposal",
            },
        )

        self.assertEqual(response.status_code, 400)
        self.assertIn("file", response.json())

    def test_100장을_초과한_PowerPoint는_업로드할_수_없다(self):
        from io import BytesIO

        from pptx import Presentation

        presentation = Presentation()
        for _index in range(101):
            presentation.slides.add_slide(presentation.slide_layouts[6])
        buffer = BytesIO()
        presentation.save(buffer)
        self.client.force_login(self.user)

        response = self.client.post(
            "/api/company-documents/",
            {
                "file": SimpleUploadedFile(
                    "too-long.pptx",
                    buffer.getvalue(),
                    content_type=(
                        "application/vnd.openxmlformats-officedocument."
                        "presentationml.presentation"
                    ),
                ),
                "document_type": "proposal",
            },
        )

        self.assertEqual(response.status_code, 400)
        self.assertIn("100장", str(response.json()))

    def test_회사_문서는_최대_10개까지만_저장한다(self):
        self.client.force_login(self.user)
        for index in range(10):
            CompanyDocument.objects.create(
                user=self.user,
                file=f"company_documents/test-{index}.pdf",
                original_name=f"test-{index}.pdf",
                document_type=CompanyDocument.DocumentType.PROPOSAL,
            )

        response = self.client.post(
            "/api/company-documents/",
            {
                "file": SimpleUploadedFile("extra.pdf", b"extra"),
                "document_type": "proposal",
            },
        )

        self.assertEqual(response.status_code, 400)
        self.assertEqual(response.json()["error"], "회사 문서는 최대 10개까지 업로드할 수 있습니다.")

    def test_허용하지_않는_파일은_업로드할_수_없다(self):
        self.client.force_login(self.user)

        response = self.client.post(
            "/api/company-documents/",
            {
                "file": SimpleUploadedFile("program.exe", b"unsafe"),
                "document_type": "proposal",
            },
        )

        self.assertEqual(response.status_code, 400)
        self.assertIn("file", response.json())

    def test_본인이_업로드한_문서를_삭제한다(self):
        self.client.force_login(self.user)
        document = CompanyDocument.objects.create(
            user=self.user,
            file=SimpleUploadedFile("delete.pdf", b"delete"),
            original_name="delete.pdf",
            document_type=CompanyDocument.DocumentType.PROPOSAL,
        )

        response = self.client.delete(f"/api/company-documents/{document.id}/")

        self.assertEqual(response.status_code, 204)
        self.assertFalse(CompanyDocument.objects.filter(id=document.id).exists())


class CompanyKnowledgeTests(TestCase):
    def setUp(self):
        self.user = get_user_model().objects.create_user(
            username="knowledge-user",
            password="test-password",
        )
        self.media_directory = tempfile.TemporaryDirectory()
        self.settings_override = override_settings(
            MEDIA_ROOT=self.media_directory.name
        )
        self.settings_override.enable()
        self.addCleanup(self.settings_override.disable)
        self.addCleanup(self.media_directory.cleanup)
        self.document = CompanyDocument.objects.create(
            user=self.user,
            file=SimpleUploadedFile("company-introduction.docx", b"document"),
            original_name="company-introduction.docx",
            document_type=CompanyDocument.DocumentType.COMPANY_INTRODUCTION,
        )

    def test_회사_지식은_출처와_함께_저장된다(self):
        item = CompanyKnowledgeItem.objects.create(
            user=self.user,
            source_document=self.document,
            category=CompanyKnowledgeItem.Category.CAPABILITY,
            title="온라인 교육 운영 역량",
            content="전국 단위 온라인 교육 운영 경험을 보유하고 있습니다.",
            source_locations=["3페이지"],
        )

        self.assertEqual(item.source_locations, ["3페이지"])

    def test_원본_문서를_삭제하면_추출한_지식도_삭제된다(self):
        CompanyKnowledgeItem.objects.create(
            user=self.user,
            source_document=self.document,
            category=CompanyKnowledgeItem.Category.PERFORMANCE,
            title="교육 운영 실적",
            content="공공기관 교육을 수행했습니다.",
        )

        self.document.delete()

        self.assertEqual(CompanyKnowledgeItem.objects.count(), 0)

    @patch(
        "bids.services.company_knowledge._extract_batch_knowledge"
    )
    @patch("bids.services.company_knowledge.extract_document")
    def test_문서에서_출처가_있는_회사_지식을_추출한다(
        self,
        mock_extract,
        mock_extract_batch,
    ):
        from langchain_core.documents import Document

        from bids.services.company_knowledge import (
            ExtractedKnowledgeBatch,
            ExtractedKnowledgeItem,
            prepare_company_knowledge,
        )
        from bids.services.rag.extract_document import ExtractionResult

        mock_extract.return_value = ExtractionResult(
            documents=[
                Document(
                    page_content=(
                        "전국 15개 지점에서 공공기관 온라인 교육을 운영했습니다."
                    ),
                    metadata={"location": "3페이지"},
                )
            ],
            processed_files=["company-introduction.docx"],
        )
        mock_extract_batch.return_value = ExtractedKnowledgeBatch(
            items=[
                ExtractedKnowledgeItem(
                    category="performance",
                    title="전국 단위 교육 운영 실적",
                    content="전국 15개 지점에서 공공기관 교육을 운영했습니다.",
                    source_numbers=[1],
                    tags=["교육", "공공기관"],
                )
            ]
        )

        result = prepare_company_knowledge(self.document)
        item = CompanyKnowledgeItem.objects.get()

        self.assertFalse(result["reused"])
        self.assertEqual(result["item_count"], 1)
        self.assertEqual(item.category, "performance")
        self.assertEqual(item.source_locations, ["3페이지"])
        self.assertIn("전국 15개 지점", item.evidence_excerpt)

    @patch(
        "bids.services.company_knowledge._extract_batch_knowledge"
    )
    def test_이미_추출한_문서는_OpenAI를_다시_호출하지_않는다(
        self,
        mock_extract_batch,
    ):
        CompanyKnowledgeItem.objects.create(
            user=self.user,
            source_document=self.document,
            category=CompanyKnowledgeItem.Category.STRENGTH,
            title="전담 운영 조직",
            content="전담 운영 조직을 보유하고 있습니다.",
        )
        from bids.services.company_knowledge import prepare_company_knowledge

        result = prepare_company_knowledge(self.document)

        self.assertTrue(result["reused"])
        self.assertEqual(result["item_count"], 1)
        mock_extract_batch.assert_not_called()


class RecommendedBidViewTests(TestCase):
    def setUp(self):
        self.user = get_user_model().objects.create_user(
            username="recommend-user",
            password="test-password",
        )

    def create_profile(self, keywords):
        return CompanyProfile.objects.create(
            user=self.user,
            company_name="추천테스트회사",
            business_registration_number="111-11-11111",
            representative_name="김대표",
            address="서울특별시",
            industry="소프트웨어 개발",
            main_business="웹 서비스 개발",
            preferred_keywords=keywords,
        )

    @staticmethod
    def create_notice(index, title, **overrides):
        values = {
            "bid_ntce_no": f"REC-{index:03d}",
            "bid_ntce_ord": "000",
            "title": title,
            "is_active": True,
            "raw_data": {
                "bidNtceNo": f"REC-{index:03d}",
                "bidNtceOrd": "000",
                "bidNtceNm": title,
            },
        }
        values.update(overrides)
        return BidNotice.objects.create(**values)

    def test_로그인하지_않으면_추천_공고를_볼_수_없다(self):
        response = self.client.get("/api/recommended-bids/")

        self.assertEqual(response.status_code, 401)

    def test_회사_프로필이_없으면_빈_목록을_반환한다(self):
        self.client.force_login(self.user)

        response = self.client.get("/api/recommended-bids/")

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()["count"], 0)
        self.assertEqual(response.json()["items"], [])

    def test_회사_키워드가_포함된_유효_공고만_추천한다(self):
        self.create_profile("홈페이지, 시스템 유지보수")
        self.create_notice(1, "공공기관 홈페이지 구축 용역")
        self.create_notice(2, "업무 시스템 유지보수 사업")
        self.create_notice(3, "사무실 청소 용역")
        self.create_notice(4, "마감된 홈페이지 구축", is_active=False)
        self.client.force_login(self.user)

        response = self.client.get("/api/recommended-bids/")
        data = response.json()
        notice_numbers = {item["bidNtceNo"] for item in data["items"]}

        self.assertEqual(response.status_code, 200)
        self.assertEqual(data["keywords"], ["홈페이지", "시스템 유지보수"])
        self.assertEqual(data["count"], 2)
        self.assertEqual(notice_numbers, {"REC-001", "REC-002"})

    def test_기관명이나_허용업종에_키워드가_있어도_추천한다(self):
        self.create_profile("소프트웨어")
        self.create_notice(
            1,
            "정보시스템 사업",
            allowed_industry="소프트웨어사업자",
        )
        self.create_notice(
            2,
            "장비 구매",
            notice_organization="소프트웨어진흥원",
        )
        self.client.force_login(self.user)

        response = self.client.get("/api/recommended-bids/")

        self.assertEqual(response.json()["count"], 2)
        self.assertTrue(response.json()["items"][0]["matchReasons"])

    def test_임시_키워드로_검색해도_회사_키워드는_변경하지_않는다(self):
        profile = self.create_profile("홈페이지")
        self.create_notice(1, "공공기관 홈페이지 구축")
        self.create_notice(2, "방탄 장비 구매")
        self.create_notice(3, "교육용 장비 구매")
        self.client.force_login(self.user)

        response = self.client.get(
            "/api/recommended-bids/",
            {"keywords": "방탄, 구매"},
        )
        data = response.json()

        profile.refresh_from_db()  # DB에 저장된 회사 키워드를 다시 확인
        self.assertEqual(data["keywords"], ["방탄", "구매"])
        self.assertEqual(data["count"], 2)
        self.assertEqual(profile.preferred_keywords, "홈페이지")

    def test_회사가_선택한_공고_유형만_추천한다(self):
        profile = self.create_profile("홈페이지")
        profile.preferred_bid_type = "service"
        profile.save(update_fields=["preferred_bid_type"])
        self.create_notice(1, "홈페이지 구축 용역", business_type="용역")
        self.create_notice(2, "홈페이지 장비 구매", business_type="물품")
        self.client.force_login(self.user)

        response = self.client.get("/api/recommended-bids/")
        data = response.json()

        self.assertEqual(data["count"], 1)
        self.assertEqual(data["items"][0]["bidNtceNo"], "REC-001")

    def test_회사_희망지역과_참가_가능지역이_맞는_공고를_추천한다(self):
        profile = self.create_profile("홈페이지")
        profile.preferred_region = "서울"
        profile.save(update_fields=["preferred_region"])
        self.create_notice(1, "홈페이지 구축", region_limit=False)
        self.create_notice(2, "서울 홈페이지 구축", region_limit=True, allowed_region="서울특별시")
        self.create_notice(3, "부산 홈페이지 구축", region_limit=True, allowed_region="부산광역시")
        self.client.force_login(self.user)

        response = self.client.get("/api/recommended-bids/")
        data = response.json()

        self.assertEqual(data["region"], "서울")
        self.assertEqual(data["count"], 2)

    def test_임시_희망지역으로_검색해도_회사_지역은_변경하지_않는다(self):
        profile = self.create_profile("홈페이지")
        profile.preferred_region = "서울"
        profile.save(update_fields=["preferred_region"])
        self.create_notice(1, "부산 홈페이지 구축", region_limit=True, allowed_region="부산광역시")
        self.create_notice(2, "서울 홈페이지 구축", region_limit=True, allowed_region="서울특별시")
        self.client.force_login(self.user)

        response = self.client.get(
            "/api/recommended-bids/",
            {"region": "부산"},
        )

        profile.refresh_from_db()
        self.assertEqual(response.json()["count"], 1)
        self.assertEqual(response.json()["items"][0]["bidNtceNo"], "REC-001")
        self.assertEqual(profile.preferred_region, "서울")

    def test_추천_공고는_최대_20건만_반환한다(self):
        self.create_profile("홈페이지")
        for index in range(25):
            self.create_notice(index, f"홈페이지 구축 사업 {index}")
        self.client.force_login(self.user)

        response = self.client.get("/api/recommended-bids/")
        data = response.json()

        self.assertEqual(data["count"], 25)
        self.assertEqual(len(data["items"]), 20)


class DocumentExtractionTests(TestCase):
    def test_HWPX_텍스트와_문서위치를_추출한다(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = Path(temp_dir) / "request.hwpx"
            with zipfile.ZipFile(file_path, "w") as archive:
                archive.writestr(
                    "Contents/section0.xml",
                    "<root xmlns:hp='urn:test'><hp:t>참가 자격</hp:t></root>",
                )

            result = extract_document(file_path)

        self.assertEqual(result.failed_files, [])
        self.assertIn("참가 자격", result.documents[0].page_content)
        self.assertEqual(result.documents[0].metadata["location"], "문서 구역 1")

    def test_ZIP의_문서는_처리하고_실행파일은_실패로_기록한다(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = Path(temp_dir) / "attachments.zip"
            with zipfile.ZipFile(file_path, "w") as archive:
                archive.writestr("notice.txt", "입찰 공고 내용")
                archive.writestr("danger.exe", b"MZ")

            result = extract_document(file_path)

        self.assertEqual(result.processed_files, ["notice.txt"])
        self.assertEqual(result.failed_files[0]["file_name"], "danger.exe")
        self.assertIn("실행 가능한 파일", result.failed_files[0]["reason"])


class BidDocumentSearchTests(TestCase):
    def test_Chroma_관련도_점수를_0과_1_사이로_제한한다(self):
        from bids.services.rag.vector_store import normalize_l2_relevance_score

        self.assertEqual(normalize_l2_relevance_score(0), 1)
        self.assertEqual(normalize_l2_relevance_score(10), 0)
        self.assertGreater(normalize_l2_relevance_score(0.5), 0)

    @patch("bids.services.rag.retriever.get_bid_vector_store")
    def test_관련도_기준을_통과한_chunk를_모두_선택한다(self, mock_get_store):
        documents = [object() for _ in range(10)]
        vector_store = mock_get_store.return_value
        vector_store._collection.count.return_value = 10
        vector_store.similarity_search_with_relevance_scores.return_value = [
            (document, 0.9) for document in documents
        ]

        result = search_bid_documents("BID-001", "사업비가 얼마야?")

        self.assertEqual(result, documents)
        vector_store.similarity_search_with_relevance_scores.assert_called_once_with(
            "사업비가 얼마야?",
            k=10,
        )

    @patch("bids.services.rag.retriever.get_bid_vector_store")
    def test_관련_chunk가_너무_적으면_상위_5개를_선택한다(self, mock_get_store):
        documents = [object() for _ in range(8)]
        vector_store = mock_get_store.return_value
        vector_store._collection.count.return_value = 8
        vector_store.similarity_search_with_relevance_scores.return_value = [
            (document, 0.1) for document in documents
        ]

        result = search_bid_documents("BID-001", "참가 자격은?")

        self.assertEqual(result, documents[:5])


class BidChatViewTests(TestCase):
    def setUp(self):
        self.user = get_user_model().objects.create_user(
            username="chat-user",
            password="test-password",
        )
        self.notice = BidNotice.objects.create(
            bid_ntce_no="CHAT-001",
            bid_ntce_ord="000",
            title="AI 챗봇 테스트 공고",
            is_active=True,
        )

    def test_로그인하지_않으면_챗봇을_사용할_수_없다(self):
        response = self.client.post(
            "/api/bids/CHAT-001/chat/",
            {"question": "참가 자격은?"},
            content_type="application/json",
        )

        self.assertEqual(response.status_code, 401)

    def test_빈_질문은_거부한다(self):
        self.client.force_login(self.user)
        response = self.client.post(
            "/api/bids/CHAT-001/chat/",
            {"question": ""},
            content_type="application/json",
        )

        self.assertEqual(response.status_code, 400)

    def test_없는_공고번호는_404를_반환한다(self):
        self.client.force_login(self.user)
        response = self.client.post(
            "/api/bids/NOT-FOUND/chat/",
            {"question": "참가 자격은?"},
            content_type="application/json",
        )

        self.assertEqual(response.status_code, 404)

    @patch("bids.views.generate_bid_chat_answer")
    def test_AI_답변과_출처를_JSON으로_반환한다(self, mock_generate_answer):
        mock_generate_answer.return_value = {
            "answer": "테스트 AI 답변",
            "sources": [{"file_name": "공고문.pdf", "page": "2"}],
        }
        self.client.force_login(self.user)

        response = self.client.post(
            "/api/bids/CHAT-001/chat/",
            {"question": " 참가 자격은? "},
            content_type="application/json",
        )

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()["answer"], "테스트 AI 답변")
        self.assertEqual(response.json()["sources"][0]["page"], "2")
        self.assertEqual(BidChatMessage.objects.count(), 2)
        mock_generate_answer.assert_called_once_with("CHAT-001", "참가 자격은?")

    @patch("bids.views.generate_bid_chat_answer")
    def test_저장된_채팅을_다시_불러온다(self, mock_generate_answer):
        mock_generate_answer.return_value = {
            "answer": "저장된 답변",
            "sources": [],
        }
        self.client.force_login(self.user)
        self.client.post(
            "/api/bids/CHAT-001/chat/",
            {"question": "사업비는?"},
            content_type="application/json",
        )

        response = self.client.get("/api/bids/CHAT-001/chat/")

        self.assertEqual(response.status_code, 200)
        self.assertEqual(len(response.json()["messages"]), 2)
        self.assertEqual(response.json()["messages"][1]["text"], "저장된 답변")

    @patch("bids.views.generate_bid_chat_answer")
    def test_공고별_질문은_기존_횟수와_관계없이_사용한다(self, mock_generate_answer):
        mock_generate_answer.return_value = {
            "answer": "추가 답변",
            "sources": [{"file_name": "공고문.pdf", "page": "2"}],
        }
        saved_bid = SavedBid.objects.create(
            user=self.user,
            bid_notice=self.notice,
        )
        BidChatMessage.objects.bulk_create(
            [
                BidChatMessage(
                    saved_bid=saved_bid,
                    role=BidChatMessage.Role.USER,
                    content=f"질문 {index}",
                )
                for index in range(20)
            ]
        )
        self.client.force_login(self.user)

        response = self.client.post(
            "/api/bids/CHAT-001/chat/",
            {"question": "추가 질문"},
            content_type="application/json",
        )

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()["answer"], "추가 답변")
        mock_generate_answer.assert_called_once_with("CHAT-001", "추가 질문")

    def test_proposal_revision_message_is_saved_as_pending(self):
        saved_bid = SavedBid.objects.create(
            user=self.user,
            bid_notice=self.notice,
        )
        BidProposal.objects.create(
            saved_bid=saved_bid,
            output_format=BidProposal.OutputFormat.PPTX,
            template_mode="default_template",
            strategy={},
            revision_plan={"version": "template_generation_v1"},
        )
        self.client.force_login(self.user)

        response = self.client.post(
            "/api/bids/CHAT-001/chat/",
            {"message": "8페이지 수행 전략을 수정해줘"},
            content_type="application/json",
        )

        message = BidChatMessage.objects.get(role=BidChatMessage.Role.USER)
        assistant_message = BidChatMessage.objects.get(
            role=BidChatMessage.Role.ASSISTANT,
        )
        self.assertEqual(response.status_code, 201)
        self.assertEqual(response.json()["intent"], "revision")
        self.assertEqual(
            message.message_type,
            BidChatMessage.MessageType.PROPOSAL,
        )
        self.assertEqual(message.status, BidChatMessage.Status.PENDING)
        self.assertIn("반영할까요", assistant_message.content)

    def test_chat_message_delete_keeps_deleted_placeholder(self):
        saved_bid = SavedBid.objects.create(
            user=self.user,
            bid_notice=self.notice,
        )
        message = BidChatMessage.objects.create(
            saved_bid=saved_bid,
            role=BidChatMessage.Role.USER,
            content="삭제할 메시지",
        )
        self.client.force_login(self.user)

        response = self.client.delete(
            f"/api/bids/CHAT-001/chat/{message.id}/"
        )

        message.refresh_from_db()
        self.assertEqual(response.status_code, 200)
        self.assertTrue(message.is_deleted)
        self.assertEqual(message.content, "")
        self.assertEqual(
            response.json()["message"]["text"],
            "삭제된 메시지입니다.",
        )


class BidListViewTests(TestCase):
    @staticmethod
    def create_notice(index, **overrides):
        values = {
            "bid_ntce_no": f"BID-{index:03d}",
            "bid_ntce_ord": "000",
            "title": f"Test notice {index}",
            "deadline_status": BidNotice.DeadlineStatus.ACTIVE,
            "is_active": True,
            "raw_data": {
                "bidNtceNo": f"BID-{index:03d}",
                "bidNtceOrd": "000",
                "bidNtceNm": f"Test notice {index}",
            },
        }
        values.update(overrides)
        return BidNotice.objects.create(**values)

    def test_returns_twenty_active_notices_per_page(self):
        for index in range(25):
            self.create_notice(index)

        self.create_notice(
            25,
            deadline_status=BidNotice.DeadlineStatus.REVIEW,
        )
        self.create_notice(
            26,
            deadline_status=BidNotice.DeadlineStatus.EXPIRED,
            is_active=False,
        )

        response = self.client.get("/api/bids/")
        data = response.json()

        self.assertEqual(response.status_code, 200)
        self.assertEqual(data["count"], 26)
        self.assertEqual(data["page"], 1)
        self.assertEqual(data["page_size"], 20)
        self.assertEqual(data["total_pages"], 2)
        self.assertEqual(len(data["items"]), 20)

    def test_공고번호로_상세공고_한건을_조회한다(self):
        self.create_notice(
            1,
            title="제주 정보시스템 구축 용역",
            region_limit=True,
            allowed_region="제주특별자치도",
        )

        response = self.client.get("/api/bids/BID-001/")  # 상세 API 주소 요청
        item = response.json()["item"]

        self.assertEqual(response.status_code, 200)
        self.assertEqual(item["bidNtceNo"], "BID-001")
        self.assertEqual(item["allowedRegion"], "제주특별자치도")

    def test_없는_공고번호는_404를_반환한다(self):
        response = self.client.get("/api/bids/NOT-FOUND/")

        self.assertEqual(response.status_code, 404)
        self.assertIn("error", response.json())

    def test_returns_second_page_and_deadline_status(self):
        for index in range(21):
            self.create_notice(index)

        response = self.client.get("/api/bids/?page=2&page_size=20")
        data = response.json()

        self.assertEqual(response.status_code, 200)
        self.assertEqual(len(data["items"]), 1)
        self.assertEqual(data["items"][0]["deadlineStatus"], "active")
        self.assertTrue(data["items"][0]["isActive"])

    def test_참가가능지역과_마지막업데이트시간을_반환한다(self):
        self.create_notice(
            1,
            region_limit=True,
            allowed_region="제주특별자치도",
        )

        response = self.client.get("/api/bids/")
        item = response.json()["items"][0]

        self.assertTrue(item["regionLimit"])
        self.assertEqual(item["allowedRegion"], "제주특별자치도")
        self.assertIsNotNone(response.json()["last_updated_at"])

    def test_rejects_invalid_page_parameters(self):
        response = self.client.get("/api/bids/?page=wrong&page_size=20")

        self.assertEqual(response.status_code, 400)
        self.assertIn("page", response.json()["error"])

    def test_searches_notice_title_and_organization(self):
        self.create_notice(
            1,
            title="School repair project",
            notice_organization="Seoul Office of Education",
        )
        self.create_notice(2, title="Hospital supplies")

        response = self.client.get("/api/bids/?q=School")

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()["count"], 1)
        self.assertEqual(response.json()["items"][0]["bidNtceNo"], "BID-001")

    def test_회사_키워드와_여러_희망지역으로_공고를_필터한다(self):
        self.create_notice(
            1,
            title="공공기관 홈페이지 구축",
            region_limit=True,
            allowed_region="충청북도",
        )
        self.create_notice(
            2,
            title="정보시스템 유지보수",
            region_limit=True,
            allowed_region="제주특별자치도",
        )
        self.create_notice(3, title="사무용품 구매", region_limit=False)

        response = self.client.get(
            "/api/bids/",
            {"keywords": "홈페이지,시스템", "regions": "충북,제주"},
        )

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()["count"], 2)

    def test_filters_business_type_deadline_status_and_region(self):
        self.create_notice(
            1,
            business_type="공사",
            deadline_status=BidNotice.DeadlineStatus.REVIEW,
            region_limit=True,
            allowed_region="서울특별시",
        )
        self.create_notice(
            2,
            business_type="용역",
            deadline_status=BidNotice.DeadlineStatus.REVIEW,
            region_limit=True,
            allowed_region="서울특별시",
        )
        self.create_notice(3, business_type="공사")

        response = self.client.get(
            "/api/bids/?business_type=공사&deadline_status=review&region=seoul"
        )

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()["count"], 1)
        self.assertEqual(response.json()["summary"]["construction"], 1)

    def test_참가_지역은_지역제한이_없거나_선택지역인_공고를_반환한다(self):
        self.create_notice(1, region_limit=False)
        self.create_notice(2, region_limit=True, allowed_region="서울특별시")
        self.create_notice(3, region_limit=True, allowed_region="부산광역시")

        response = self.client.get("/api/bids/?region=seoul")

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()["count"], 2)

    def test_마감_상태로_공고를_필터한다(self):
        self.create_notice(1, deadline_status=BidNotice.DeadlineStatus.ACTIVE)
        self.create_notice(2, deadline_status=BidNotice.DeadlineStatus.REVIEW)

        response = self.client.get("/api/bids/?deadline_status=review")

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()["count"], 1)
        self.assertEqual(response.json()["items"][0]["bidNtceNo"], "BID-002")

    def test_filters_notices_by_deadline_days(self):
        self.create_notice(1, close_at=timezone.now() + timedelta(days=2))
        self.create_notice(2, close_at=timezone.now() + timedelta(days=10))

        response = self.client.get("/api/bids/?deadline_days=3")

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()["count"], 1)

    def test_마감일을_빠른순과_느린순으로_정렬한다(self):
        self.create_notice(1, close_at=timezone.now() + timedelta(days=2))
        self.create_notice(2, close_at=timezone.now() + timedelta(days=10))

        early_response = self.client.get("/api/bids/?deadline_sort=asc")
        late_response = self.client.get("/api/bids/?deadline_sort=desc")

        self.assertEqual(early_response.json()["items"][0]["bidNtceNo"], "BID-001")
        self.assertEqual(late_response.json()["items"][0]["bidNtceNo"], "BID-002")

    def test_공고일을_빠른순과_느린순으로_정렬한다(self):
        self.create_notice(1, notice_date=timezone.localdate() - timedelta(days=2))
        self.create_notice(2, notice_date=timezone.localdate())

        early_response = self.client.get("/api/bids/?notice_sort=asc")
        late_response = self.client.get("/api/bids/?notice_sort=desc")

        self.assertEqual(early_response.json()["items"][0]["bidNtceNo"], "BID-001")
        self.assertEqual(late_response.json()["items"][0]["bidNtceNo"], "BID-002")

    def test_계약방법으로_공고를_필터한다(self):
        self.create_notice(1, contract_method="제한경쟁")
        self.create_notice(2, contract_method="수의계약")

        response = self.client.get("/api/bids/?contract_method=수의계약")

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()["count"], 1)
        self.assertEqual(response.json()["items"][0]["bidNtceNo"], "BID-002")


class BidSyncViewTests(TestCase):
    def setUp(self):
        self.user = get_user_model().objects.create_user(
            username="sync-user",
            password="test-password",
        )

    def test_로그인하지_않으면_공고를_업데이트할_수_없다(self):
        response = self.client.post("/api/bids/sync/")

        self.assertEqual(response.status_code, 401)

    @patch("bids.views.call_command")
    def test_공고_수집명령을_실행한다(self, mock_call_command):
        self.client.force_login(self.user)

        response = self.client.post("/api/bids/sync/")

        self.assertEqual(response.status_code, 200)
        mock_call_command.assert_called_once()


class SavedBidModelTests(TestCase):
    def setUp(self):
        self.user = get_user_model().objects.create_user(
            username="saved-bid-user",
            password="test-password",
        )
        self.notice = BidNotice.objects.create(
            bid_ntce_no="SAVE-001",
            bid_ntce_ord="000",
            title="저장 테스트 공고",
        )

    def test_회원과_공고를_연결해서_저장한다(self):
        saved_bid = SavedBid.objects.create(
            user=self.user,
            bid_notice=self.notice,
        )

        self.assertEqual(saved_bid.user, self.user)
        self.assertEqual(saved_bid.bid_notice, self.notice)

    def test_같은_회원은_같은_공고를_중복저장할_수_없다(self):
        SavedBid.objects.create(user=self.user, bid_notice=self.notice)

        with self.assertRaises(IntegrityError):
            with transaction.atomic():
                SavedBid.objects.create(user=self.user, bid_notice=self.notice)


class SavedBidViewTests(TestCase):
    def setUp(self):
        self.user = get_user_model().objects.create_user(
            username="saved-api-user",
            password="test-password",
        )
        self.other_user = get_user_model().objects.create_user(
            username="other-saved-api-user",
            password="test-password",
        )
        self.notice = BidNotice.objects.create(
            bid_ntce_no="SAVE-API-001",
            bid_ntce_ord="000",
            title="저장 API 테스트 공고",
            is_active=True,
            raw_data={
                "bidNtceNo": "SAVE-API-001",
                "bidNtceOrd": "000",
                "bidNtceNm": "저장 API 테스트 공고",
            },
        )

    def test_로그인하지_않으면_저장공고_API를_사용할_수_없다(self):
        response = self.client.get("/api/saved-bids/")

        self.assertEqual(response.status_code, 401)

    def test_공고번호로_공고를_저장한다(self):
        self.client.force_login(self.user)

        response = self.client.post(
            "/api/saved-bids/",
            {"bid_ntce_no": "SAVE-API-001"},
            content_type="application/json",
        )

        self.assertEqual(response.status_code, 201)
        self.assertTrue(response.json()["created"])
        self.assertTrue(SavedBid.objects.filter(user=self.user).exists())

    def test_같은_공고를_다시_저장해도_중복되지_않는다(self):
        self.client.force_login(self.user)

        for _ in range(2):
            response = self.client.post(
                "/api/saved-bids/",
                {"bid_ntce_no": "SAVE-API-001"},
                content_type="application/json",
            )

        self.assertEqual(response.status_code, 200)
        self.assertFalse(response.json()["created"])
        self.assertEqual(SavedBid.objects.filter(user=self.user).count(), 1)

    def test_로그인한_회원의_저장공고만_조회한다(self):
        SavedBid.objects.create(user=self.user, bid_notice=self.notice)
        SavedBid.objects.create(user=self.other_user, bid_notice=self.notice)
        self.client.force_login(self.user)

        response = self.client.get("/api/saved-bids/")

        self.assertEqual(response.json()["count"], 1)
        self.assertEqual(response.json()["items"][0]["bidNtceNo"], "SAVE-API-001")
        self.assertFalse(response.json()["items"][0]["hasChat"])
        self.assertFalse(response.json()["items"][0]["hasAnalysis"])

    def test_내가_저장한_공고만_저장취소한다(self):
        SavedBid.objects.create(user=self.user, bid_notice=self.notice)
        SavedBid.objects.create(user=self.other_user, bid_notice=self.notice)
        self.client.force_login(self.user)

        response = self.client.delete("/api/saved-bids/SAVE-API-001/")

        self.assertEqual(response.status_code, 204)
        self.assertFalse(SavedBid.objects.filter(user=self.user).exists())
        self.assertTrue(SavedBid.objects.filter(user=self.other_user).exists())


class BidSyncHelperTests(TestCase):
    def test_마감일이_미래면_유효_상태다(self):
        result = determine_deadline_status(
            "일반공고",
            timezone.now() + timedelta(days=1),
        )

        self.assertEqual(result, BidNotice.DeadlineStatus.ACTIVE)

    def test_마감일이_없으면_확인_필요_상태다(self):
        result = determine_deadline_status("일반공고", None)

        self.assertEqual(result, BidNotice.DeadlineStatus.REVIEW)

    def test_취소공고는_마감_상태다(self):
        result = determine_deadline_status("취소공고", None)

        self.assertEqual(result, BidNotice.DeadlineStatus.EXPIRED)

    def test_같은_공고번호에서는_높은_차수를_선택한다(self):
        latest_notices = {}
        update_latest_notices(
            latest_notices,
            [
                {"bidNtceNo": "R001", "bidNtceOrd": "000"},
                {"bidNtceNo": "R001", "bidNtceOrd": "002"},
                {"bidNtceNo": "R001", "bidNtceOrd": "001"},
            ],
        )

        self.assertEqual(latest_notices["R001"]["bidNtceOrd"], "002")

    @override_settings(G2B_API_KEY="test-key")
    @patch("bids.services.g2b_api.requests.get")
    def test_API에_페이지_조건을_전달한다(self, mock_get):
        mock_get.return_value.json.return_value = {"response": {"body": {}}}
        start_at = datetime(2026, 7, 1, tzinfo=ZoneInfo("Asia/Seoul"))
        end_at = datetime(2026, 7, 12, tzinfo=ZoneInfo("Asia/Seoul"))

        fetch_bid_notices(
            page_no=2,
            num_of_rows=999,
            start_at=start_at,
            end_at=end_at,
        )

        params = mock_get.call_args.kwargs["params"]
        self.assertEqual(params["pageNo"], 2)
        self.assertEqual(params["numOfRows"], 999)
        self.assertEqual(params["bidNtceBgnDt"], "202607010000")
        self.assertEqual(params["bidNtceEndDt"], "202607120000")


class RecommendationTests(TestCase):
    def setUp(self):
        self.user = get_user_model().objects.create_user(
            username="recommend-user",
            password="1234",
        )
        CompanyProfile.objects.create(
            user=self.user,
            company_name="추천 테스트 회사",
            business_registration_number="111-22-33333",
            representative_name="김대표",
            address="서울특별시",
            industry="외국어 교육",
            main_business="외국어 교육 서비스",
            preferred_keywords="외국어, 교육",
            preferred_bid_type="service",
            preferred_region="서울",
        )
        self.notice = BidNotice.objects.create(
            bid_ntce_no="REC-001",
            title="서울 외국어 교육 운영 용역",
            business_type="용역",
            region_limit=True,
            allowed_region="서울특별시",
            is_active=True,
            raw_data={
                "bidNtceNo": "REC-001",
                "bidNtceOrd": "000",
                "bidNtceNm": "서울 외국어 교육 운영 용역",
                "bsnsDivNm": "용역",
            },
        )

    def test_회사조건과_일치한_공고를_한번만_추천한다(self):
        first = match_user_recommendations(self.user)
        second = match_user_recommendations(self.user)

        self.assertEqual(first["created"], 1)
        self.assertEqual(second["created"], 0)
        self.assertEqual(RecommendedBid.objects.count(), 1)

    def test_저장된_추천공고_API는_현재_회원의_공고를_보여준다(self):
        match_user_recommendations(self.user)
        self.client.force_login(self.user)

        response = self.client.get("/api/recommendations/")

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()["count"], 1)
        self.assertEqual(response.json()["items"][0]["bidNtceNo"], "REC-001")

    def test_마감일이_지난_추천공고는_API에서_제외한다(self):
        self.notice.close_at = timezone.now() - timedelta(minutes=1)
        self.notice.save(update_fields=["close_at"])
        RecommendedBid.objects.create(
            user=self.user,
            bid_notice=self.notice,
            match_score=60,
            is_match=True,
        )
        self.client.force_login(self.user)

        response = self.client.get("/api/recommendations/")

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()["count"], 0)
        self.assertTrue(RecommendedBid.objects.filter(user=self.user).exists())

    def test_마감일이_지난_공고는_추천_매칭에서_다시_만들지_않는다(self):
        self.notice.close_at = timezone.now() - timedelta(minutes=1)
        self.notice.save(update_fields=["close_at"])

        result = match_user_recommendations(self.user)

        self.assertEqual(result["created"], 0)
        self.assertFalse(RecommendedBid.objects.filter(user=self.user).exists())

    def test_회사조건이_바뀌면_예전_추천을_목록에서_숨긴다(self):
        match_user_recommendations(self.user)
        profile = self.user.company_profile
        profile.preferred_keywords = "건축"
        profile.save(update_fields=["preferred_keywords"])

        match_user_recommendations(self.user)
        self.client.force_login(self.user)
        response = self.client.get("/api/recommendations/")

        self.assertEqual(response.json()["count"], 0)
        self.assertEqual(RecommendedBid.objects.count(), 1)

    def test_제외_키워드가_포함된_공고는_추천하지_않는다(self):
        profile = self.user.company_profile
        profile.excluded_keywords = "운영"
        profile.save(update_fields=["excluded_keywords"])

        result = match_user_recommendations(self.user)

        self.assertEqual(result["created"], 0)
        self.assertFalse(RecommendedBid.objects.filter(is_match=True).exists())

    def test_기존_필수와_관심_키워드를_하나의_목록으로_추천한다(self):
        profile = self.user.company_profile
        profile.required_keywords = "번역"
        profile.preferred_keywords = "외국어"
        profile.save(update_fields=["required_keywords", "preferred_keywords"])
        BidNotice.objects.create(
            bid_ntce_no="REC-002",
            title="서울 번역 운영 용역",
            business_type="용역",
            region_limit=True,
            allowed_region="서울특별시",
            is_active=True,
        )

        match_user_recommendations(self.user)

        self.assertEqual(
            RecommendedBid.objects.filter(user=self.user, is_match=True).count(),
            2,
        )

    def test_공고명에_일치한_키워드가_많을수록_점수가_높다(self):
        BidNotice.objects.create(
            bid_ntce_no="REC-002",
            title="서울 외국어 통역 운영 용역",
            business_type="용역",
            region_limit=True,
            allowed_region="서울특별시",
            is_active=True,
        )

        match_user_recommendations(self.user)

        two_keyword_score = RecommendedBid.objects.get(
            bid_notice__bid_ntce_no="REC-001"
        ).match_score
        one_keyword_score = RecommendedBid.objects.get(
            bid_notice__bid_ntce_no="REC-002"
        ).match_score
        self.assertEqual(two_keyword_score, 70)
        self.assertEqual(one_keyword_score, 60)
        self.assertGreater(two_keyword_score, one_keyword_score)

    def test_점수가_같으면_공고명_키워드가_많이_일치한_공고가_먼저다(self):
        BidNotice.objects.create(
            bid_ntce_no="REC-002",
            title="서울 외국어 운영 용역",
            allowed_industry="교육 서비스",
            business_type="용역",
            region_limit=True,
            allowed_region="서울특별시",
            is_active=True,
            raw_data={
                "bidNtceNo": "REC-002",
                "bidNtceOrd": "000",
                "bidNtceNm": "서울 외국어 운영 용역",
                "bsnsDivNm": "용역",
            },
        )
        match_user_recommendations(self.user)
        self.client.force_login(self.user)

        response = self.client.get("/api/recommendations/")

        self.assertEqual(response.json()["items"][0]["bidNtceNo"], "REC-001")
        self.assertEqual(
            list(
                RecommendedBid.objects.order_by("bid_notice__bid_ntce_no").values_list(
                    "match_score", "title_match_count"
                )
            ),
            [(70, 2), (70, 1)],
        )

    def test_조건에_맞는_공고가_많아도_상위_10건만_활성화한다(self):
        for number in range(2, 13):
            BidNotice.objects.create(
                bid_ntce_no=f"REC-{number:03d}",
                title=f"서울 외국어 교육 운영 용역 {number}",
                business_type="용역",
                region_limit=True,
                allowed_region="서울특별시",
                is_active=True,
            )

        match_user_recommendations(self.user)

        self.assertEqual(
            RecommendedBid.objects.filter(user=self.user, is_match=True).count(),
            10,
        )


class BidAnalysisTests(TestCase):
    def setUp(self):
        self.user = get_user_model().objects.create_user(
            username="analysis-user",
            password="1234",
        )
        self.other_user = get_user_model().objects.create_user(
            username="other-analysis-user",
            password="1234",
        )
        CompanyProfile.objects.create(
            user=self.user,
            company_name="분석 테스트 회사",
            business_registration_number="222-33-44444",
            representative_name="이대표",
            address="서울특별시",
            industry="소프트웨어",
            main_business="시스템 개발",
            preferred_keywords="시스템",
        )
        self.notice = BidNotice.objects.create(
            bid_ntce_no="ANALYSIS-001",
            title="정보시스템 구축 용역",
            notice_organization="테스트 기관",
            business_type="용역",
            is_active=True,
            raw_data={
                "bidNtceNo": "ANALYSIS-001",
                "bidNtceOrd": "000",
                "bidNtceNm": "정보시스템 구축 용역",
                "bsnsDivNm": "용역",
            },
        )
        self.saved_bid = SavedBid.objects.create(
            user=self.user,
            bid_notice=self.notice,
        )
        self.report = {
            "summary": "회사 역량과 공고 분야가 대체로 일치합니다.",
            "fit_score": 72,
            "recommendation": "조건부 검토",
            "overview": {
                "ordering_organization": "테스트 기관",
                "budget": "1억원",
                "bid_deadline": "2026-07-31",
                "contract_period": "3개월",
                "project_summary": "정보시스템 구축",
            },
            "evaluation_items": [
                {
                    "name": "업종/분야 적합성",
                    "score": 15,
                    "max_score": 20,
                    "status": "충족",
                    "explanation": "주요 사업과 일치",
                    "source_numbers": [1],
                }
            ],
            "eligibility": ["소프트웨어사업자"],
            "required_documents": ["사업자등록증"],
            "technical_evaluation": ["기술 제안 평가"],
            "price_evaluation": ["가격 점수 평가"],
            "main_tasks": ["시스템 구축"],
            "required_staff": ["PM 1명"],
            "certifications_and_experience": ["유사 실적"],
            "contract_cautions": ["납기 확인"],
            "strengths": ["개발 역량"],
            "risks": ["실적 확인 필요"],
            "company_checks": ["실적증명서 확인"],
            "action_strategy": ["담당자 배정"],
            "sources": [
                {"number": 1, "file_name": "공고문.pdf", "location": "3페이지"}
            ],
        }

    @patch("bids.services.rag.analysis.generate_bid_analysis")
    def test_AI분석은_한번_생성한_결과를_재사용한다(self, mock_generate):
        mock_generate.return_value = self.report
        self.client.force_login(self.user)

        first = self.client.post("/api/bids/ANALYSIS-001/analysis/")
        second = self.client.post("/api/bids/ANALYSIS-001/analysis/")

        self.assertEqual(first.status_code, 201)
        self.assertEqual(second.status_code, 200)
        self.assertEqual(mock_generate.call_count, 1)
        self.assertEqual(BidAnalysis.objects.count(), 1)

    def test_다른_회원은_내_공고_분석에_접근할_수_없다(self):
        self.client.force_login(self.other_user)

        response = self.client.get("/api/bids/ANALYSIS-001/analysis/")

        self.assertEqual(response.status_code, 404)

    def test_저장된_분석을_PDF로_받는다(self):
        BidAnalysis.objects.create(saved_bid=self.saved_bid, report=self.report)
        self.client.force_login(self.user)

        response = self.client.get("/api/bids/ANALYSIS-001/analysis/pdf/")
        content = b"".join(response.streaming_content)

        self.assertEqual(response.status_code, 200)
        self.assertTrue(content.startswith(b"%PDF"))


class BidProposalTests(TestCase):
    def setUp(self):
        self.media_directory = tempfile.TemporaryDirectory()
        self.media_override = override_settings(MEDIA_ROOT=self.media_directory.name)
        self.media_override.enable()
        self.addCleanup(self.media_override.disable)
        self.addCleanup(self.media_directory.cleanup)

        self.user = get_user_model().objects.create_user(
            username="proposal-user",
            password="1234",
        )
        CompanyProfile.objects.create(
            user=self.user,
            company_name="제안서 테스트 회사",
            business_registration_number="333-44-55555",
            representative_name="박대표",
            address="서울특별시",
            industry="소프트웨어",
            main_business="정보시스템 구축",
            preferred_keywords="정보시스템",
        )
        self.notice = BidNotice.objects.create(
            bid_ntce_no="PROPOSAL-001",
            title="정보시스템 구축 제안요청",
            business_type="용역",
            is_active=True,
            raw_data={
                "bidNtceNo": "PROPOSAL-001",
                "bidNtceOrd": "000",
                "bidNtceNm": "정보시스템 구축 제안요청",
                "bsnsDivNm": "용역",
            },
        )
        self.saved_bid = SavedBid.objects.create(
            user=self.user,
            bid_notice=self.notice,
        )

    @patch("bids.services.proposal_preview.get_template_slide_preview")
    def test_제안서_템플릿_슬라이드_이미지를_조회한다(self, mock_preview):
        image_path = Path(self.media_directory.name) / "template-slide.png"
        image_path.write_bytes(b"template-preview")
        mock_preview.return_value = (image_path, 30)
        self.client.force_login(self.user)

        response = self.client.get(
            "/api/proposal-templates/corporate/slides/1/"
        )
        content = b"".join(response.streaming_content)

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response["X-Slide-Count"], "30")
        self.assertEqual(content, b"template-preview")

    def create_source_proposal(self):
        from io import BytesIO

        from pptx import Presentation

        presentation = Presentation()
        first_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        first_slide.shapes.title.text = "2022년 한국가스공사 제안"
        first_slide.placeholders[1].text = "기존 사업 수행 전략"
        second_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        second_slide.shapes.title.text = "이전 사업 일정"
        second_slide.placeholders[1].text = "2022년 1월부터 6월까지 수행"
        buffer = BytesIO()
        presentation.save(buffer)

        return CompanyDocument.objects.create(
            user=self.user,
            file=SimpleUploadedFile(
                "previous-proposal.pptx",
                buffer.getvalue(),
                content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ),
            original_name="previous-proposal.pptx",
            document_type=CompanyDocument.DocumentType.PROPOSAL,
        )

    def test_제안서_만들기를_누르면_프로젝트가_저장된다(self):
        self.client.force_login(self.user)

        response = self.client.post(
            "/api/saved-bids/PROPOSAL-001/proposal-project/"
        )

        self.assertEqual(response.status_code, 200)
        self.saved_bid.refresh_from_db()
        self.assertIsNotNone(self.saved_bid.proposal_started_at)
        self.assertTrue(response.json()["hasProposalProject"])

    def test_프로젝트만_삭제하고_저장공고는_유지한다(self):
        self.client.force_login(self.user)
        self.saved_bid.proposal_started_at = timezone.now()
        self.saved_bid.save(update_fields=["proposal_started_at"])
        BidChatMessage.objects.create(
            saved_bid=self.saved_bid,
            role=BidChatMessage.Role.USER,
            content="제안서 수정 요청",
        )

        response = self.client.delete(
            "/api/saved-bids/PROPOSAL-001/proposal-project/"
        )

        self.assertEqual(response.status_code, 204)
        self.saved_bid.refresh_from_db()
        self.assertIsNone(self.saved_bid.proposal_started_at)
        self.assertTrue(SavedBid.objects.filter(id=self.saved_bid.id).exists())
        self.assertFalse(self.saved_bid.chat_messages.exists())

    def test_제안서에_공고기본정보를_전달한다(self):
        from bids.services.rag.proposal import build_bid_notice_context

        context = build_bid_notice_context(self.notice)

        self.assertIn("정보시스템 구축 제안요청", context)
        self.assertIn("PROPOSAL-001", context)
        self.assertIn("용역", context)

    @patch("bids.services.rag.proposal.get_bid_vector_store")
    @patch("bids.services.rag.proposal.search_bid_documents")
    def test_제안서에는_관련_chunk와_나머지_공고문서도_함께_전달한다(
        self,
        mock_search,
        mock_vector_store,
    ):
        from langchain_core.documents import Document

        from bids.services.rag.proposal import collect_proposal_documents

        relevant = Document(
            page_content="핵심 평가 기준",
            metadata={
                "source": "request.pdf",
                "file_name": "request.pdf",
                "element_index": 1,
                "location": "1페이지",
            },
        )
        mock_search.return_value = [relevant]
        mock_vector_store.return_value.get.return_value = {
            "documents": [
                "핵심 평가 기준",
                "세부 과업 내용",
                "계약 주의사항",
            ],
            "metadatas": [
                relevant.metadata,
                {
                    "source": "request.pdf",
                    "file_name": "request.pdf",
                    "element_index": 2,
                    "location": "2페이지",
                },
                {
                    "source": "notice.hwp",
                    "file_name": "notice.hwp",
                    "element_index": 1,
                    "location": "1문단",
                },
            ],
        }

        documents = collect_proposal_documents("PROPOSAL-001")

        self.assertEqual(documents[0].page_content, "핵심 평가 기준")
        self.assertEqual(
            {document.page_content for document in documents},
            {"핵심 평가 기준", "세부 과업 내용", "계약 주의사항"},
        )

    def test_PPTX_슬라이드와_텍스트_위치를_읽는다(self):
        from bids.services.proposal_pptx_renderer import extract_pptx_inventory

        source_proposal = self.create_source_proposal()
        inventory = extract_pptx_inventory(source_proposal.file.path)

        self.assertEqual(len(inventory), 2)
        self.assertEqual(inventory[0]["title"], "2022년 한국가스공사 제안")
        self.assertEqual(inventory[0]["elements"][0]["target"], "shape-0")
        self.assertEqual(inventory[0]["elements"][0]["kind"], "title")
        self.assertGreater(inventory[0]["elements"][0]["max_chars"], 0)

    def test_묶음별_작성에도_전체_슬라이드_구성을_전달한다(self):
        from bids.services.rag.proposal import build_deck_outline

        outline = build_deck_outline(
            [
                {"slide_number": 1, "title": "표지", "role": "cover"},
                {"slide_number": 2, "title": "수행 전략", "role": "content"},
            ]
        )

        self.assertIn("1. 표지 (역할: cover)", outline)
        self.assertIn("2. 수행 전략 (역할: content)", outline)

    def test_생성본의_자리표시자와_빈_슬라이드를_자동_검수한다(self):
        from io import BytesIO

        from pptx import Presentation

        from bids.services.proposal_pptx_renderer import inspect_proposal_quality

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "[사업명]"
        slide.placeholders[1].text = "짧은 문구"
        buffer = BytesIO()
        presentation.save(buffer)

        result = inspect_proposal_quality(buffer.getvalue())

        self.assertFalse(result["passed"])
        self.assertEqual(
            result["unresolved_placeholders"][0]["marker"],
            "[사업명]",
        )
        self.assertEqual(result["empty_slide_numbers"], [1])

    def test_제안서_작성규칙에_업종별_목차가_저장되어있다(self):
        from bids.services.proposal_rules import (
            build_proposal_rules_context,
            load_proposal_rules,
        )

        rules = load_proposal_rules()
        context = build_proposal_rules_context()

        self.assertEqual(rules["reference_summary"]["reviewed_slide_count"], 448)
        self.assertIn("education_service", rules["domain_modules"])
        self.assertIn("발주처 요구", context)

    def test_텍스트상자_수용량을_크게_넘는_수정은_건너뛴다(self):
        from io import BytesIO

        from pptx import Presentation
        from pptx.util import Inches, Pt

        from bids.services.proposal_pptx_renderer import build_proposal_pptx

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = slide.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(1.2),
            Inches(0.35),
        )
        textbox.text_frame.paragraphs[0].text = "기존 문구"
        textbox.text_frame.paragraphs[0].runs[0].font.size = Pt(20)
        source = BytesIO()
        presentation.save(source)

        with tempfile.TemporaryDirectory() as directory:
            source_path = Path(directory) / "source.pptx"
            source_path.write_bytes(source.getvalue())
            result = build_proposal_pptx(
                source_path=source_path,
                bid_notice=self.notice,
                revision_plan={
                    "slide_changes": [
                        {
                            "slide_number": 1,
                            "action": "UPDATE",
                            "title": "긴 문구 테스트",
                            "reason": "겹침 방지 확인",
                            "text_changes": [
                                {
                                    "target": "shape-0",
                                    "original_text": "기존 문구",
                                    "revised_text": "가" * 500,
                                    "reason": "긴 문구",
                                }
                            ],
                        }
                    ],
                    "added_slides": [],
                },
            )

        generated = Presentation(BytesIO(result["file_bytes"]))
        self.assertEqual(generated.slides[0].shapes[0].text, "기존 문구")
        self.assertTrue(result["revision_log"][0]["warnings"])

    def test_100장_제안서를_25장씩_4개_묶음으로_나눈다(self):
        from bids.services.rag.proposal import split_slide_inventory

        inventory = [
            {
                "slide_number": number,
                "title": f"{number}페이지",
                "elements": [],
                "role": "content",
            }
            for number in range(1, 101)
        ]

        batches = split_slide_inventory(inventory)

        self.assertEqual([len(batch) for batch in batches], [25, 25, 25, 25])
        self.assertEqual(batches[0][0]["slide_number"], 1)
        self.assertEqual(batches[-1][-1]["slide_number"], 100)

    def test_묶음별_개정계획을_전체_제안서_규칙으로_통합한다(self):
        from bids.services.rag.proposal import (
            merge_revision_batch_plans,
            split_slide_inventory,
        )

        inventory = [
            {
                "slide_number": number,
                "title": f"{number}페이지",
                "elements": [],
                "role": "content",
            }
            for number in range(1, 101)
        ]
        batch_reviews = []
        for batch in split_slide_inventory(inventory):
            start_slide = batch[0]["slide_number"]
            batch_reviews.append(
                {
                    "inventory": batch,
                    "plan": {
                        "summary": f"{start_slide}페이지부터 검토",
                        "slide_changes": [
                            {
                                "slide_number": start_slide,
                                "action": "REMOVE",
                                "title": f"{start_slide}페이지",
                                "reason": "현재 공고와 무관",
                                "text_changes": [],
                                "source_numbers": [],
                            },
                            {
                                "slide_number": start_slide + 1,
                                "action": "REMOVE",
                                "title": f"{start_slide + 1}페이지",
                                "reason": "현재 공고와 무관",
                                "text_changes": [],
                                "source_numbers": [],
                            },
                        ],
                        "added_slides": [
                            {
                                "after_slide_number": start_slide,
                                "template_slide_number": start_slide + 1,
                                "title": f"추가 전략 {start_slide}-{index}",
                                "reason": "공고 전략 보강",
                                "text_changes": [],
                                "source_numbers": [],
                            }
                            for index in range(6)
                        ],
                        "final_review_items": [f"{start_slide}페이지 확인"],
                    },
                }
            )

        merged = merge_revision_batch_plans(batch_reviews)

        self.assertEqual(len(merged["review_batches"]), 4)
        self.assertEqual(
            sum(
                batch["reviewed_slide_count"]
                for batch in merged["review_batches"]
            ),
            100,
        )
        self.assertEqual(
            [
                change["action"]
                for change in merged["slide_changes"]
            ].count("REMOVE"),
            3,
        )
        self.assertEqual(len(merged["added_slides"]), 20)

    @patch("bids.services.rag.proposal.build_revision_chain")
    def test_100장_제안서는_AI_검토를_4번_순차_실행한다(
        self,
        mock_build_revision_chain,
    ):
        from bids.services.rag.proposal import build_revision_plan_in_batches

        class FakeRevisionResult:
            def model_dump(self):
                return {
                    "summary": "묶음 검토 완료",
                    "slide_changes": [],
                    "added_slides": [],
                    "final_review_items": [],
                }

        inventory = [
            {
                "slide_number": number,
                "title": f"{number}페이지",
                "elements": [],
                "role": "content",
            }
            for number in range(1, 101)
        ]
        revision_chain = mock_build_revision_chain.return_value
        revision_chain.invoke.side_effect = [
            FakeRevisionResult()
            for _ in range(4)
        ]

        result = build_revision_plan_in_batches(
            inventory,
            {
                "company_context": "회사 정보",
                "bid_notice_context": "공고 기본정보",
                "bid_context": "공고 문서",
                "strategy_context": "수주 전략",
                "source_proposal_context": "기존 제안서",
                "allowed_template_numbers": "2, 3",
            },
        )

        self.assertEqual(revision_chain.invoke.call_count, 4)
        self.assertIn(
            "1~25페이지",
            revision_chain.invoke.call_args_list[0].args[0]["batch_scope"],
        )
        self.assertIn(
            "76~100페이지",
            revision_chain.invoke.call_args_list[-1].args[0]["batch_scope"],
        )
        self.assertEqual(len(result["review_batches"]), 4)

    def test_원본_PPTX를_수정_삭제_추가하고_메모를_남긴다(self):
        from io import BytesIO

        from pptx import Presentation

        from bids.services.proposal_pptx_renderer import build_proposal_pptx

        source_proposal = self.create_source_proposal()
        result = build_proposal_pptx(
            source_path=source_proposal.file.path,
            bid_notice=self.notice,
            revision_plan={
                "slide_changes": [
                    {
                        "slide_number": 1,
                        "action": "UPDATE",
                        "title": "표지",
                        "reason": "새 사업명 반영",
                        "text_changes": [
                            {
                                "target": "shape-0",
                                "original_text": "2022년 한국가스공사 제안",
                                "revised_text": "정보시스템 구축 제안",
                                "reason": "사업명 변경",
                            }
                        ],
                    },
                    {
                        "slide_number": 2,
                        "action": "REMOVE",
                        "title": "이전 사업 일정",
                        "reason": "현재 공고와 무관",
                        "text_changes": [],
                    },
                ],
                "added_slides": [
                    {
                        "after_slide_number": 1,
                        "template_slide_number": 2,
                        "title": "신규 수행 전략",
                        "reason": "공고 요구사항 보강",
                        "text_changes": [
                            {
                                "target": "shape-0",
                                "content_label": "슬라이드 제목",
                                "original_text": "이전 사업 일정",
                                "revised_text": "신규 수행 전략",
                                "reason": "추가 슬라이드 제목",
                            },
                            {
                                "target": "shape-1",
                                "content_label": "수행 내용",
                                "original_text": "2022년 1월부터 6월까지 수행",
                                "revised_text": "새 공고 일정에 맞춘 수행 계획",
                                "reason": "현재 사업 일정 반영",
                            }
                        ],
                    }
                ],
            },
        )

        generated = Presentation(BytesIO(result["file_bytes"]))
        self.assertEqual(len(generated.slides), 2)
        self.assertEqual(generated.slides[0].shapes.title.text, "정보시스템 구축 제안")
        self.assertEqual(generated.slides[1].shapes.title.text, "신규 수행 전략")
        self.assertIn("AI 제안서 개정 메모", generated.slides[0].notes_slide.notes_text_frame.text)
        self.assertIn(
            "UPDATE",
            [entry["action"] for entry in result["revision_log"]],
        )

    def test_표지_복제와_대량_슬라이드_삭제를_차단한다(self):
        from io import BytesIO

        from pptx import Presentation

        from bids.services.proposal_pptx_renderer import build_proposal_pptx

        presentation = Presentation()
        for index in range(1, 7):
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = f"{index}페이지 제목"
            slide.placeholders[1].text = f"{index}페이지 본문"
        source = BytesIO()
        presentation.save(source)

        with tempfile.TemporaryDirectory() as directory:
            source_path = Path(directory) / "source.pptx"
            source_path.write_bytes(source.getvalue())
            result = build_proposal_pptx(
                source_path=source_path,
                bid_notice=self.notice,
                revision_plan={
                    "slide_changes": [
                        {
                            "slide_number": index,
                            "action": "REMOVE",
                            "title": f"{index}페이지",
                            "reason": "삭제 요청",
                            "text_changes": [],
                        }
                        for index in range(1, 7)
                    ],
                    "added_slides": [
                        {
                            "after_slide_number": 1,
                            "template_slide_number": 1,
                            "title": "표지 복제",
                            "reason": "잘못된 추가 요청",
                            "text_changes": [
                                {
                                    "target": "shape-0",
                                    "original_text": "1페이지 제목",
                                    "revised_text": "복제 제목",
                                    "reason": "제목 수정",
                                },
                                {
                                    "target": "shape-1",
                                    "original_text": "1페이지 본문",
                                    "revised_text": "복제 본문",
                                    "reason": "본문 수정",
                                },
                            ],
                        }
                    ],
                },
            )

        generated = Presentation(BytesIO(result["file_bytes"]))
        self.assertEqual(len(generated.slides), 3)
        self.assertEqual(generated.slides[0].shapes.title.text, "1페이지 제목")
        self.assertNotIn(
            "ADD",
            [entry["action"] for entry in result["revision_log"]],
        )
        self.assertEqual(
            [entry["action"] for entry in result["revision_log"]].count("REMOVE"),
            3,
        )

    def test_같은_본문_디자인으로_최대_20장을_추가한다(self):
        from io import BytesIO

        from pptx import Presentation

        from bids.services.proposal_pptx_renderer import build_proposal_pptx

        source_proposal = self.create_source_proposal()
        additions = []
        for index in range(1, 22):
            additions.append(
                {
                    "after_slide_number": 2,
                    "template_slide_number": 2,
                    "title": f"추가 전략 {index}",
                    "reason": "공고 맞춤 전략 보강",
                    "text_changes": [
                        {
                            "target": "shape-0",
                            "original_text": "이전 사업 일정",
                            "revised_text": f"추가 전략 {index}",
                            "reason": "제목 변경",
                        },
                        {
                            "target": "shape-1",
                            "original_text": "2022년 1월부터 6월까지 수행",
                            "revised_text": f"새 공고 수행 방안 {index}",
                            "reason": "본문 변경",
                        },
                    ],
                }
            )

        result = build_proposal_pptx(
            source_path=source_proposal.file.path,
            bid_notice=self.notice,
            revision_plan={
                "slide_changes": [],
                "added_slides": additions,
            },
        )

        generated = Presentation(BytesIO(result["file_bytes"]))
        self.assertEqual(len(generated.slides), 22)
        self.assertEqual(
            [entry["action"] for entry in result["revision_log"]].count("ADD"),
            20,
        )

    @patch("bids.services.company_knowledge.prepare_user_company_knowledge")
    def test_제안서에_자동_추출한_회사_지식을_전달한다(self, mock_prepare):
        from bids.services.company_knowledge import build_company_knowledge_context

        source_document = CompanyDocument.objects.create(
            user=self.user,
            file=SimpleUploadedFile("company-introduction.pptx", b"document"),
            original_name="company-introduction.pptx",
            document_type=CompanyDocument.DocumentType.COMPANY_INTRODUCTION,
        )
        CompanyKnowledgeItem.objects.create(
            user=self.user,
            source_document=source_document,
            category=CompanyKnowledgeItem.Category.CAPABILITY,
            title="공공 정보시스템 운영 역량",
            content="공공 정보시스템 구축과 운영 경험을 보유하고 있습니다.",
            source_locations=["1페이지"],
        )
        mock_prepare.return_value = {
            "item_count": 1,
            "processed_files": [],
            "reused_files": ["company-introduction.pptx"],
            "failed_files": [],
        }

        context, processing = build_company_knowledge_context(self.user)

        self.assertIn("공공 정보시스템 구축과 운영 경험", context)
        self.assertEqual(processing["item_count"], 1)

    def test_제안서에_사업자등록_기본정보를_전달한다(self):
        from bids.services.rag.analysis import company_context

        profile = CompanyProfile.objects.get(user=self.user)
        context = company_context(profile)

        self.assertIn("333-44-55555", context)
        self.assertIn("박대표", context)

    @patch("bids.services.rag.proposal.create_bid_proposal_from_template")
    def test_Bid2_템플릿으로_새_제안서를_생성한다(self, mock_create):
        mock_create.return_value = {
            "strategy": {"win_themes": ["공고 맞춤 전략"]},
            "revision_plan": {
                "version": "template_generation_v1",
                "summary": "Bid2 템플릿으로 작성했습니다.",
                "revision_log": [],
            },
            "template_mode": "default_template",
            "filename": "proposal-PROPOSAL-001.pptx",
            "file_bytes": b"generated proposal",
            "content_type": (
                "application/vnd.openxmlformats-officedocument."
                "presentationml.presentation"
            ),
        }
        self.client.force_login(self.user)

        response = self.client.post(
            "/api/bids/PROPOSAL-001/proposal/",
            {
                "generation_mode": "default_template",
                "template_id": "public",
            },
            content_type="application/json",
        )

        proposal = BidProposal.objects.get(saved_bid=self.saved_bid)
        self.assertEqual(response.status_code, 201)
        self.assertEqual(proposal.template_mode, "default_template")
        self.assertEqual(proposal.revision_plan["template_id"], "public")
        mock_create.assert_called_once()

    def test_선택할_수_있는_제안서_템플릿을_조회한다(self):
        self.client.force_login(self.user)

        response = self.client.get("/api/bids/PROPOSAL-001/proposal/")

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()["selected_template_id"], "public")
        self.assertEqual(
            [item["id"] for item in response.json()["templates"]],
            ["corporate", "modern", "public"],
        )

    def test_제안서를_삭제하면_PDF_미리보기_캐시도_삭제한다(self):
        from bids.services.proposal_preview import get_proposal_preview_path

        proposal = BidProposal(
            saved_bid=self.saved_bid,
            output_format=BidProposal.OutputFormat.PPTX,
            template_mode="default_template",
        )
        proposal.generated_file.save(
            "proposal.pptx",
            SimpleUploadedFile("proposal.pptx", b"proposal"),
        )
        preview_path = get_proposal_preview_path(proposal)
        preview_path.write_bytes(b"%PDF-preview")

        proposal.delete()

        self.assertFalse(preview_path.exists())

    def test_legacy_source_revision_is_not_exposed_as_current_proposal(self):
        BidProposal.objects.create(
            saved_bid=self.saved_bid,
            output_format=BidProposal.OutputFormat.PPTX,
            template_mode="source_revision",
            strategy={},
            revision_plan={
                "version": "pptx_revision_v3",
                "output_slide_count": 98,
            },
        )
        self.client.force_login(self.user)

        proposal_response = self.client.get(
            "/api/bids/PROPOSAL-001/proposal/"
        )
        saved_bids_response = self.client.get("/api/saved-bids/")

        self.assertEqual(proposal_response.status_code, 200)
        self.assertIsNone(proposal_response.json()["proposal"])
        self.assertFalse(saved_bids_response.json()["items"][0]["hasProposal"])

    def test_기존_PPT_수정_방식은_사용할_수_없다(self):
        self.client.force_login(self.user)

        response = self.client.post(
            "/api/bids/PROPOSAL-001/proposal/",
            {"generation_mode": "existing", "source_document_id": 1},
            content_type="application/json",
        )

        self.assertEqual(response.status_code, 400)
        self.assertIn("등록된 제안서 템플릿", response.json()["error"])

    def test_저장된_제안서_파일을_내려받는다(self):
        proposal = BidProposal.objects.create(
            saved_bid=self.saved_bid,
            output_format=BidProposal.OutputFormat.PPTX,
            strategy={},
            revision_plan={},
        )
        proposal.generated_file.save(
            "revised-proposal-PROPOSAL-001.pptx",
            SimpleUploadedFile("proposal.pptx", b"generated proposal"),
        )
        self.client.force_login(self.user)

        response = self.client.get(
            "/api/bids/PROPOSAL-001/proposal/download/"
        )
        content = b"".join(response.streaming_content)

        self.assertEqual(response.status_code, 200)
        self.assertEqual(content, b"generated proposal")

    @patch("bids.services.rag.proposal.revise_proposal_with_feedback")
    def test_미리보기_수정요청을_초안에_저장한다(self, mock_feedback):
        proposal = BidProposal.objects.create(
            saved_bid=self.saved_bid,
            output_format=BidProposal.OutputFormat.PPTX,
            template_mode="default_template",
            strategy={"win_themes": ["안정적인 구축"]},
            revision_plan={
                "version": "template_generation_v1",
                "status": "draft",
                "revision_log": [],
                "feedback_history": [],
                "source_slide_count": 2,
                "output_slide_count": 2,
            },
        )
        proposal.generated_file.save(
            "draft.pptx",
            SimpleUploadedFile("draft.pptx", b"draft"),
        )
        mock_feedback.return_value = {
            "revision_plan": {
                "summary": "8페이지 수행 전략을 구체화했습니다.",
                "final_review_items": [],
            },
            "filename": "feedback.pptx",
            "file_bytes": b"feedback",
            "output_slide_count": 2,
            "revision_log": [
                {
                    "source_slide_number": 8,
                    "output_slide_number": 8,
                    "action": "UPDATE",
                    "title": "수행 전략",
                    "reason": "사용자 요청 반영",
                    "changes": [],
                    "warnings": [],
                }
            ],
        }
        self.client.force_login(self.user)

        response = self.client.post(
            "/api/bids/PROPOSAL-001/proposal/feedback/",
            {
                "instruction": "8페이지 수행 전략을 더 구체적으로 수정해 주세요.",
                "slide_number": 8,
            },
            content_type="application/json",
        )

        proposal.refresh_from_db()
        self.assertEqual(response.status_code, 200)
        self.assertEqual(proposal.revision_plan["status"], "draft")
        self.assertEqual(len(proposal.revision_plan["feedback_history"]), 1)
        self.assertEqual(
            proposal.revision_plan["feedback_history"][0]["slide_number"],
            8,
        )

    def test_초안은_확정한_뒤에만_내려받는다(self):
        proposal = BidProposal.objects.create(
            saved_bid=self.saved_bid,
            output_format=BidProposal.OutputFormat.PPTX,
            template_mode="default_template",
            strategy={},
            revision_plan={
                "version": "template_generation_v1",
                "status": "draft",
            },
        )
        proposal.generated_file.save(
            "draft.pptx",
            SimpleUploadedFile("draft.pptx", b"draft proposal"),
        )
        self.client.force_login(self.user)

        blocked = self.client.get(
            "/api/bids/PROPOSAL-001/proposal/download/"
        )
        finalized = self.client.post(
            "/api/bids/PROPOSAL-001/proposal/finalize/"
        )
        downloaded = self.client.get(
            "/api/bids/PROPOSAL-001/proposal/download/"
        )
        content = b"".join(downloaded.streaming_content)

        self.assertEqual(blocked.status_code, 409)
        self.assertEqual(finalized.status_code, 200)
        self.assertEqual(finalized.json()["proposal"]["status"], "final")
        self.assertEqual(downloaded.status_code, 200)
        self.assertEqual(content, b"draft proposal")
